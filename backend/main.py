from __future__ import annotations
"""
SynaptDI — FastAPI Backend  v2.0
Endpoints:
  POST /query                      → RAG answer
  GET  /documents                  → list uploaded documents
  POST /documents/upload           → upload + background-ingest a file
  GET  /documents/{file}/status    → ingestion progress
  DELETE /documents/{file}         → remove from ChromaDB + disk
  GET  /branding                   → branding config
  POST /branding                   → save branding config
  GET  /users                      → user list
  POST /users                      → add user
  PUT  /users/{id}                 → update user
  DELETE /users/{id}               → remove user
  GET  /stats/queries              → per-doc query hit counts
  GET  /health  GET /stats         → health / index stats
"""

import os, json, time, hashlib, uuid, re, requests, threading, base64, hmac, secrets, shutil, subprocess
from pathlib import Path
from typing import Optional

from fastapi import FastAPI, HTTPException, UploadFile, File, BackgroundTasks, Header, Depends
from fastapi.responses import StreamingResponse, FileResponse
import mimetypes
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import chromadb
import ingest   # reuse the ingestion parsers: parse_openapi_spec, parse_markdown, chunk_text, embed_batch, …
import conformance   # TMF630 conformance rule engine
import tmf_profile   # profile-aware conformance (diff vs canonical TMF specs)
import xray          # API estate X-ray (portfolio roll-up)
import oda           # ODA Component conformance (schema roll-up)
import oda_component_contract  # canonical ODA Component contract resolver (execution-backed conformance)
import oda_ctk_jobs   # ODA Component CTK execution jobs (execution-backed conformance)
import oda_ctk_sanitize   # centralized public-response redaction boundary (Phase 4 audit)
import spec_facts    # deterministic answers for structured spec-fact questions
import knowledge_entities   # Knowledge Engine V2 — deterministic entity extraction (Phase 7B)
import knowledge_link       # Knowledge Engine V2 — corpus inventory + entity linking
import knowledge_router     # Knowledge Engine V2 — deterministic routing before RAG
from vectorstore import get_store, reset_store   # backend-agnostic vector store (Chroma today)
import vectorstore   # Phase 7C: read the active/default knowledge collection for startup reporting

_chroma_lock = threading.Lock()   # protects singleton + upsert/query overlap

# ── Config ─────────────────────────────────────────────────────────────────────
OLLAMA_URL  = os.getenv("OLLAMA_URL",  "http://localhost:11434")
LLM_MODEL   = os.getenv("LLM_MODEL",  "llama3.1:8b")   # PRD wants an 8B; 3.3 is 70B-only, so 3.1:8b is the real 8B
FAST_MODEL  = os.getenv("FAST_MODEL", "llama3.2:latest")  # 3B — "fast mode" for quick lookups
# Serialize LLM generations so concurrent users queue instead of thrashing Ollama.
LLM_CONCURRENCY = int(os.getenv("LLM_CONCURRENCY", "1"))
# ── Performance knobs (matter a lot on CPU-only machines) ───────────────────────
NUM_CTX          = int(os.getenv("NUM_CTX", "4096"))          # prompt+gen window
KEEP_ALIVE       = os.getenv("OLLAMA_KEEP_ALIVE", "30m")      # keep model resident between queries (no cold reload)
DEEP_NUM_PREDICT = int(os.getenv("DEEP_NUM_PREDICT", "900"))  # Deep-mode max answer tokens
FAST_NUM_PREDICT = int(os.getenv("FAST_NUM_PREDICT", "500"))  # Fast-mode: shorter answer = quicker
FAST_TOP_K       = int(os.getenv("FAST_TOP_K", "4"))          # Fast-mode: fewer chunks = faster prompt processing
NUM_THREAD       = int(os.getenv("OLLAMA_NUM_THREAD", "0"))   # 0 = let Ollama decide; set to physical cores on CPU
EMBED_MODEL = os.getenv("EMBED_MODEL","nomic-embed-text")
CHROMA_PATH = os.getenv("CHROMA_PATH","./chroma_db")
# Response cache: on by default; set CACHE_ENABLED=false to force every /query and
# /query/stream call through the real, current pipeline (no reads, no writes) — the
# safe way to smoke-test a routing/grounding/integrity change against live traffic
# without a stale cached answer masking it.
CACHE_ENABLED = os.getenv("CACHE_ENABLED", "1").strip().lower() not in ("0", "false", "no", "off")
# Bump whenever routing, grounding, integrity-validation, or version-resolution
# semantics change — it's folded into the cache key, so a cached answer from an older
# pipeline is never reused (a plain cache_clear() catches "content changed"; this
# catches "the code that turns content into an answer changed"). History: 1 = original
# deterministic short-circuit (raw template, no LLM); 2 = grounded-generation bridge +
# integrity validation; 3 = required-fields router rewrite (confirm/challenge phrasing,
# explicit-version isolation, independent-version comparison); 4 = integrity validator's
# "not required" hedge-detection widened to cover contractions/adverbs/articles
# ("isn't required", "not actually required", "not a required field") — the narrower
# version was misreading a correct denial as a false claim and discarding good answers.
CACHE_PIPELINE_VERSION = os.getenv("CACHE_PIPELINE_VERSION", "5")  # Phase 7B: Knowledge Engine V2 (routing/version/abstention) — invalidates all V1 cached answers

# ── Retrieval tuning ─────────────────────────────────────────────────────────
TOP_K            = 8      # chunks handed to the LLM
KB_DIST_MAX      = 0.85   # cosine-distance ceiling for a KB chunk to count as relevant
CONFIDENT_DIST   = 0.35   # if the best chunk is this close, treat as high-confidence in-domain
                          # and use the no-refusal prompt even when no spec ID was named
                          # (stops the model over-refusing on strong but un-named matches)
# Uploaded docs are intentional context — treat them as first-class, not a gated afterthought.
UPLOAD_MAX_DIST  = 0.70   # include an uploaded chunk if it's at least loosely on-topic
UPLOAD_BOOST     = 0.06   # let a relevant upload edge out an equidistant KB chunk in ranking
UPLOAD_CONFIDENT = 0.55   # an uploaded chunk this close → answer confidently from the user's doc
DOC_SCOPE_MAX    = 0.90   # in "My Documents" scope, accept almost anything from the docs

NO_INFO_MSG    = ("I don't have enough information in the current knowledge base to "
                  "answer this. The relevant document may not be indexed yet.")
NO_INFO_PREFIX = "i don't have enough information in the current knowledge base"

STORAGE_DIR = Path("./storage"); STORAGE_DIR.mkdir(exist_ok=True)
UPLOADS_DIR = Path("./uploads");  UPLOADS_DIR.mkdir(exist_ok=True)

UPLOADS_FILE  = STORAGE_DIR / "uploads.json"
BRANDING_FILE = STORAGE_DIR / "branding.json"
USERS_FILE    = STORAGE_DIR / "users.json"
CHATCFG_FILE  = STORAGE_DIR / "chat_config.json"
CONVOS_DIR    = STORAGE_DIR / "conversations"; CONVOS_DIR.mkdir(exist_ok=True)
PROJECTS_DIR  = STORAGE_DIR / "projects";      PROJECTS_DIR.mkdir(exist_ok=True)
ANALYTICS_FILE = STORAGE_DIR / "analytics.json"
FEEDBACK_FILE  = STORAGE_DIR / "feedback.json"
_analytics_lock = threading.Lock()

def _log_query(question: str, answered: bool, user: str = ""):
    """Track query volume, per-day series, per-user activity, and 'knowledge gaps'
    (questions we couldn't answer). Powers the adoption dashboard."""
    try:
        with _analytics_lock:
            a = json.loads(ANALYTICS_FILE.read_text()) if ANALYTICS_FILE.exists() else {}
            a["total"] = a.get("total", 0) + 1
            if answered:
                a["answered"] = a.get("answered", 0) + 1
            else:
                gaps = a.get("gaps", [])
                gaps.insert(0, {"q": question[:200], "ts": int(time.time())})
                a["gaps"] = gaps[:50]
            day = time.strftime("%Y-%m-%d")
            days = a.get("days", {})
            d = days.get(day, {"q": 0, "a": 0, "users": {}})
            d["q"] += 1
            if answered:
                d["a"] += 1
            who = (user or "anonymous").lower()
            d["users"][who] = d["users"].get(who, 0) + 1
            days[day] = d
            a["days"] = {k: days[k] for k in sorted(days)[-90:]}   # keep 90 days
            ANALYTICS_FILE.write_text(json.dumps(a))
    except Exception:
        pass

# ── In-memory state ────────────────────────────────────────────────────────────
PROCESSING_STATUS: dict[str, dict] = {}   # file → {status, chunks, error?}
QUERY_HITS:        dict[str, int]  = {}   # source → hit count (resets on restart)

# Persisted, dated per-document performance: {date: {doc: {hits, rel_sum}}}.
DOC_PERF_FILE = STORAGE_DIR / "doc_perf.json"
_doc_perf_lock = threading.Lock()

def _record_doc_hits(hits: list):
    """hits: [(doc_name, distance|None)] retrieved for one query. Persists a real,
    dated hit count + summed relevance (1 - cosine distance) per document, so the
    Document Performance tab shows genuine per-period data instead of estimates."""
    if not hits:
        return
    try:
        with _doc_perf_lock:
            data = json.loads(DOC_PERF_FILE.read_text()) if DOC_PERF_FILE.exists() else {}
            day = time.strftime("%Y-%m-%d")
            bucket = data.get(day, {})
            for name, dist in hits:
                rel = max(0.0, min(1.0, 1.0 - float(dist))) if isinstance(dist, (int, float)) else 0.0
                e = bucket.get(name) or {"hits": 0, "rel_sum": 0.0}
                e["hits"] += 1
                e["rel_sum"] += rel
                bucket[name] = e
            data[day] = bucket
            DOC_PERF_FILE.write_text(json.dumps({k: data[k] for k in sorted(data)[-90:]}))
    except Exception:
        pass
_llm_sem = threading.BoundedSemaphore(max(1, LLM_CONCURRENCY))

# ── Answer cache ───────────────────────────────────────────────────────────────
# Repeated questions (demos, FAQs) come back instantly instead of re-generating.
# Only confident, sourced answers are cached; cleared whenever the index changes.
CACHE_FILE = STORAGE_DIR / "answer_cache.json"
_cache_lock = threading.Lock()
_answer_cache: Optional[dict] = None

def _cache() -> dict:
    global _answer_cache
    if _answer_cache is None:
        try:    _answer_cache = json.loads(CACHE_FILE.read_text()) if CACHE_FILE.exists() else {}
        except Exception: _answer_cache = {}
    return _answer_cache

def _cache_key(question: str, scope: str, mode: str, project_instructions: str = "") -> str:
    """Every semantically-relevant input to the answer, so two questions that could
    legitimately get different answers never collide on one cache entry:
    - CACHE_PIPELINE_VERSION: bumped whenever routing/grounding/integrity/version-
      resolution logic changes — makes an old cached answer unreachable the moment the
      code that could have gotten it wrong changes, without touching stored entries.
    - question/scope/mode: as before.
    - model: the actual model name for `mode` (not just the "fast"/"deep" label) —
      changing LLM_MODEL/FAST_MODEL changes answer semantics even if mode doesn't.
    - persona: the admin-configurable domain persona baked into every system prompt
      (_system_prompt) — an admin changing it should not surface an answer generated
      under the old one.
    - project_instructions: a Project's standing notes, folded into the prompt
      (_project_extra) — two Projects asking the same question can get different,
      equally valid answers and must not share a cache entry.
    """
    model = FAST_MODEL if mode == "fast" else LLM_MODEL
    persona = (load_chat_config().get("persona") or DEFAULT_CHAT_CONFIG["persona"]).strip()
    parts = [str(CACHE_PIPELINE_VERSION), question.strip().lower(), scope, mode, model,
             persona, (project_instructions or "").strip()]
    return hashlib.sha1("|".join(parts).encode()).hexdigest()

def cache_get(question: str, scope: str, mode: str, project_instructions: str = ""):
    if not CACHE_ENABLED:
        return None
    return _cache().get(_cache_key(question, scope, mode, project_instructions))

def cache_put(question: str, scope: str, mode: str, answer: str, sources: list,
             confidence: Optional[dict] = None, project_instructions: str = ""):
    if not CACHE_ENABLED or not sources:
        return                                   # cache disabled, or a refusal/unsourced answer
    try:
        with _cache_lock:
            c = _cache()
            c[_cache_key(question, scope, mode, project_instructions)] = {
                "answer": answer, "sources": sources,
                "confidence": confidence or {}, "ts": int(time.time())}
            while len(c) > 300:                  # cap — drop oldest
                c.pop(min(c, key=lambda k: c[k].get("ts", 0)))
            CACHE_FILE.write_text(json.dumps(c))
    except Exception:
        pass

def cache_clear():
    """Drop every cached /query answer (answer_cache.json only) — called whenever the
    index content changes, and safe to call any time as a manual, development-safe
    reset (see POST /admin/cache/clear). Never touches ChromaDB, indexed documents, or
    conversation history — those live in separate storage this never opens."""
    global _answer_cache
    try:
        with _cache_lock:
            _answer_cache = {}
            CACHE_FILE.write_text("{}")
    except Exception:
        pass

# ── Default seed data ──────────────────────────────────────────────────────────
DEFAULT_BRANDING = {
    "companyName":  "SynaptDI",
    "tagline":      "Enterprise domains at your fingertips",
    "primaryColor": "#dc2626",
}

# Admin-configurable chat surface (persona feeds the system prompt; placeholder +
# suggestions drive the chat UI). Lets the same app be retargeted to any domain.
DEFAULT_CHAT_CONFIG = {
    "placeholder": "Ask about TM Forum APIs, ODA, eTOM, SID...",
    "persona":     "TM Forum telecom standards — Open APIs TMF620–TMF915, ODA (Open Digital Architecture), eTOM, and SID",
    "suggestions": [
        "What mandatory fields does TMF622 Product Order require?",
        "What is the difference between TMF620 and TMF633?",
        "How do I handle pagination in TM Forum Open APIs?",
        "Which ODA component handles trouble tickets?",
        "How do I create a service order with TMF641?",
        "What is the SID ABE for customer data?",
    ],
}

def load_chat_config() -> dict:
    if CHATCFG_FILE.exists():
        try:    return {**DEFAULT_CHAT_CONFIG, **json.loads(CHATCFG_FILE.read_text())}
        except: pass
    return dict(DEFAULT_CHAT_CONFIG)

DEFAULT_USERS = [
    {"id":"1","name":"Admin User",    "email":"admin@synaptdi.com",   "role":"admin",  "status":"active",  "lastActive":"Now",    "title":"System Administrator",      "department":"IT Operations"},
    {"id":"2","name":"Sarah Chen",    "email":"analyst@synaptdi.com", "role":"analyst","status":"active",  "lastActive":"2h ago", "title":"Telecom Standards Analyst", "department":"Architecture"},
    {"id":"3","name":"Marcus Johnson","email":"marcus@synaptdi.com",  "role":"analyst","status":"active",  "lastActive":"1d ago", "title":"Network Standards Engineer","department":"Architecture"},
    {"id":"4","name":"Lisa Park",     "email":"lisa@synaptdi.com",    "role":"viewer", "status":"away",    "lastActive":"3d ago", "title":"Product Manager",           "department":"Product"},
    {"id":"5","name":"Tom Wilson",    "email":"tom@synaptdi.com",     "role":"viewer", "status":"inactive","lastActive":"2w ago", "title":"Business Analyst",          "department":"Strategy"},
]

# Seed passwords for the default accounts — hashed on first load. Change in production.
SEED_PASSWORDS = {
    "admin@synaptdi.com":   "admin123",
    "analyst@synaptdi.com": "analyst123",
    "marcus@synaptdi.com":  "analyst123",
    "lisa@synaptdi.com":    "viewer123",
    "tom@synaptdi.com":     "viewer123",
}

# ── Auth: hashed passwords + HMAC-signed session tokens (no external deps) ────
AUTH_SECRET_FILE = STORAGE_DIR / "auth_secret"
def _auth_secret() -> bytes:
    if AUTH_SECRET_FILE.exists():
        return AUTH_SECRET_FILE.read_bytes()
    s = secrets.token_bytes(32)
    AUTH_SECRET_FILE.write_bytes(s)
    return s
AUTH_SECRET = _auth_secret()

def hash_password(password: str, salt: str = "") -> tuple:
    salt = salt or secrets.token_hex(16)
    h = hashlib.pbkdf2_hmac("sha256", password.encode(), bytes.fromhex(salt), 100_000).hex()
    return h, salt

def verify_password(password: str, h: str, salt: str) -> bool:
    if not h or not salt:
        return False
    calc, _ = hash_password(password, salt)
    return hmac.compare_digest(calc, h)

def make_token(user: dict, days: int = 7) -> str:
    payload = {"uid": user["id"], "email": user.get("email"), "role": user.get("role"),
               "exp": int(time.time()) + days * 86400}
    body = base64.urlsafe_b64encode(json.dumps(payload).encode()).decode().rstrip("=")
    sig  = hmac.new(AUTH_SECRET, body.encode(), hashlib.sha256).hexdigest()
    return f"{body}.{sig}"

def verify_token(token: str):
    try:
        body, sig = token.split(".", 1)
        expected = hmac.new(AUTH_SECRET, body.encode(), hashlib.sha256).hexdigest()
        if not hmac.compare_digest(sig, expected):
            return None
        payload = json.loads(base64.urlsafe_b64decode(body + "=" * (-len(body) % 4)))
        if payload.get("exp", 0) < time.time():
            return None
        return payload
    except Exception:
        return None

def _public_user(u: dict) -> dict:
    """User fields safe to return to clients (never the password hash/salt)."""
    return {k: u.get(k) for k in
            ("id", "name", "email", "role", "title", "department", "status", "lastActive")}

class LoginReq(BaseModel):
    email: str
    password: str

def current_user(authorization: Optional[str] = Header(None)) -> dict:
    if not authorization or not authorization.lower().startswith("bearer "):
        raise HTTPException(401, "Not authenticated")
    payload = verify_token(authorization.split(" ", 1)[1].strip())
    if not payload:
        raise HTTPException(401, "Invalid or expired session")
    u = next((x for x in _load_users() if x.get("id") == payload.get("uid")), None)
    if not u:
        raise HTTPException(401, "User no longer exists")
    return u

def require_admin(user: dict = Depends(current_user)) -> dict:
    if user.get("role") != "admin":
        raise HTTPException(403, "Admin access required")
    return user

def optional_user_email(authorization: Optional[str]) -> str:
    """Best-effort identity for analytics — '' when not logged in, never raises."""
    try:
        if not authorization or not authorization.lower().startswith("bearer "):
            return ""
        payload = verify_token(authorization.split(" ", 1)[1].strip())
        if not payload:
            return ""
        u = next((x for x in _load_users() if x.get("id") == payload.get("uid")), None)
        return (u or {}).get("email", "") or ""
    except Exception:
        return ""

# ── App ────────────────────────────────────────────────────────────────────────
app = FastAPI(
    title="SynaptDI API",
    version="2.0.0",
    description=(
        "Deterministic **TM Forum conformance** & knowledge engine — validated against "
        "**169 official TMF specs** (avg 99.9/100). Score any OpenAPI / Swagger spec against "
        "TMF630, diff it against the canonical API, complete it, and validate ODA Components — "
        "with **no LLM in the scoring**."
    ),
    openapi_tags=[
        {"name": "Conformance", "description": "Score, profile, fix, scaffold and X-ray specs against the TM Forum standard."},
        {"name": "ODA", "description": "ODA Component conformance and the official component map."},
        {"name": "Knowledge", "description": "Grounded TM Forum Q&A over the indexed spec corpus."},
    ],
)
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

def _warm_models():
    """Load the embed + LLM models into Ollama at boot so the FIRST user query
    doesn't pay the cold-start cost (which on a CPU-only box can exceed the
    client timeout). Best-effort, runs off-thread, never blocks startup."""
    try:
        embed("warmup")
    except Exception:
        pass
    try:
        requests.post(f"{OLLAMA_URL}/api/generate",
                      json={"model": LLM_MODEL, "prompt": "ok", "stream": False,
                            "keep_alive": KEEP_ALIVE, "options": {"num_predict": 1}},
                      timeout=600)
        print(f"[startup] warmed {LLM_MODEL}", flush=True)
    except Exception:
        pass

@app.on_event("startup")
def startup_init():
    # The collection is opened lazily on the first request and cached as a
    # singleton. ChromaDB's PersistentClient loads the persisted HNSW index on
    # open, so the very first query already sees the full corpus (verified).
    print("[startup] SynaptDI ready — ChromaDB opens on first request", flush=True)
    # Phase 7C: report the ACTIVE knowledge collection explicitly. axiom_v2 is the product
    # default; an explicit VECTOR_COLLECTION override is reported as such. Never silently
    # pretend V2 is active when the collection is empty/missing.
    try:
        _n = get_collection().count()
        _src = "explicit VECTOR_COLLECTION override" if vectorstore.COLLECTION_IS_EXPLICIT else "default"
        print(f"[startup] Knowledge collection: {vectorstore.COLLECTION} ({_n} chunks) [{_src}]", flush=True)
        if _n == 0:
            print(f"[startup] DEGRADED: collection '{vectorstore.COLLECTION}' is EMPTY — "
                  f"knowledge retrieval will return no evidence. Build it (ingest_v2.py for "
                  f"axiom_v2) or set VECTOR_COLLECTION to a populated collection. NOT falling "
                  f"back silently.", flush=True)
    except Exception as e:
        print(f"[startup] DEGRADED: could not open knowledge collection "
              f"'{vectorstore.COLLECTION}': {e}", flush=True)
    if os.getenv("WARM_MODELS", "1") != "0":
        threading.Thread(target=_warm_models, daemon=True).start()
    threading.Thread(target=_scheduler_loop, daemon=True).start()   # scheduled KB auto-refresh
    tmf_profile.warm_index()                                        # build the profile index off the request path

# ── Vector store ─────────────────────────────────────────────────────────────────
# Pre-built mapping of "TMF622" → "TMF622-ProductOrdering-v4.0.0.swagger.json"
# populated at startup so the query path never needs to scan all metadatas.

def get_collection():
    """Return the vector store — a thin, backend-agnostic interface (ChromaDB today,
    swappable to Milvus/Qdrant/pgvector via VECTOR_BACKEND without touching call sites)."""
    return get_store()

# ── Spec-ID awareness ────────────────────────────────────────────────────────
# Engineers name a spec directly ("mandatory fields in TMF641?"). nomic-embed-text
# clusters every telco spec into one tight region, so pure vector search often
# ranks sibling specs above the one the user actually named. We detect the spec
# ID and retrieve that spec's own chunks explicitly so they're always in context.
_SPEC_RE  = re.compile(r"TMF[\s_-]?(\d{3})")
_VER_RE   = re.compile(r"\bv?\d+(?:\.\d+){1,3}\b")   # strip "4.0.0", "v4.1.0", "2.0"
_spec_map:   Optional[dict] = None   # spec_id -> set(filenames), built once from metadata
_spec_names: list           = []     # (name_lower, spec_id), longest-first — by-name lookup

def _clean_spec_name(source: str) -> str:
    s = _VER_RE.sub(" ", source or "")
    s = re.sub(r"\bapi\b", " ", s, flags=re.I)
    s = re.sub(r"[^a-z0-9 ]", " ", s.lower())
    return re.sub(r"\s+", " ", s).strip()

def detect_spec_ids(text: str) -> list[str]:
    """Resolve specs the user named — by number (TMF641) OR by name
    ("Product Catalog Management API" → TMF620). Name matching uses the
    by-name index built in get_spec_map()."""
    t   = text.lower()
    ids = {f"TMF{n}" for n in _SPEC_RE.findall(text.upper())}
    for name, sid in _spec_names:           # specific multi-word titles only
        if name in t:
            ids.add(sid)
    return sorted(ids)

def get_spec_map(col) -> dict:
    """Map every indexed spec ID to its filename(s) AND build a by-name index
    (spec title -> ID). Built once and cached; falls back to deriving the ID
    from the filename for older indexes that predate the spec_id metadata."""
    global _spec_map, _spec_names
    if _spec_map is None:
        with _chroma_lock:
            if _spec_map is None:
                sm:    dict[str, set] = {}
                names: dict[str, str] = {}
                try:    metas = [r["metadata"] for r in col.scan()]
                except Exception: metas = []
                for m in metas:
                    sid = m.get("spec_id") or ""
                    if not sid:
                        mm  = _SPEC_RE.search((m.get("file") or "").upper())
                        sid = f"TMF{mm.group(1)}" if mm else ""
                    if not sid:
                        continue
                    sm.setdefault(sid, set()).add(m.get("file", ""))
                    nm = _clean_spec_name(m.get("source", ""))
                    # keep only specific titles (>=2 words, >=12 chars) so common
                    # words ("event", "party") can't false-match a question
                    if nm and len(nm) >= 12 and " " in nm:
                        names.setdefault(nm, sid)
                _spec_map   = sm
                _spec_names = sorted(names.items(), key=lambda kv: -len(kv[0]))
    return _spec_map

# ── Text helpers ───────────────────────────────────────────────────────────────
CHUNK_SIZE, CHUNK_OVERLAP = 800, 120           # used by the original ingest.py

# Uploads use SMALLER chunks so each chunk covers one coherent topic.
# Larger chunks produce "blurry" embeddings (average of many topics) that
# don't match specific sub-topic queries well.
UPLOAD_CHUNK_SIZE, UPLOAD_CHUNK_OVERLAP = 350, 50

def chunk_text(text: str, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    result, start = [], 0
    while start < len(text):
        result.append(text[start : start + size])
        start += size - overlap
    return [c.strip() for c in result if len(c.strip()) >= 40]

def embed(text: str) -> list[float]:
    r = requests.post(f"{OLLAMA_URL}/api/embeddings",
                      json={"model": EMBED_MODEL, "prompt": text}, timeout=30)
    r.raise_for_status()
    return r.json()["embedding"]

from functools import lru_cache
@lru_cache(maxsize=512)
def _embed_query_cached(text: str) -> tuple:
    return tuple(embed(text))

def embed_query(text: str) -> list:
    """Embedding for a retrieval query, memoised — repeat/Regenerate/same question across
    users skips the Ollama embed round-trip. Query embeddings are deterministic per model."""
    return list(_embed_query_cached(text))

def make_chunk_id(file_name: str, idx: int) -> str:
    return "up_" + hashlib.md5(f"{file_name}_{idx}".encode()).hexdigest()[:14]

# ── File text extraction ───────────────────────────────────────────────────────
def extract_text(content: bytes, file_name: str) -> str:
    name = file_name.lower()
    from io import BytesIO
    if name.endswith(".pdf"):
        import fitz
        doc = fitz.open(stream=content, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)
    if name.endswith(".docx"):
        import docx
        d = docx.Document(BytesIO(content))
        parts = [p.text for p in d.paragraphs if p.text.strip()]
        for tbl in d.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    if name.endswith(".xlsx"):
        import openpyxl
        wb = openpyxl.load_workbook(BytesIO(content), read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    if name.endswith(".pptx"):
        import pptx
        prs = pptx.Presentation(BytesIO(content))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        if any(cells):
                            parts.append(" | ".join(cells))
        return "\n".join(parts)
    return content.decode("utf-8", errors="ignore")

# ── Uploads JSON helpers ───────────────────────────────────────────────────────
def load_uploads() -> list[dict]:
    if UPLOADS_FILE.exists():
        try:    return json.loads(UPLOADS_FILE.read_text())
        except: pass
    return []

def save_uploads(docs: list[dict]):
    UPLOADS_FILE.write_text(json.dumps(docs, indent=2))

def upsert_upload(entry: dict):
    prev = next((d for d in load_uploads() if d.get("file") == entry["file"]), None)
    if prev and entry.get("folder") is None and prev.get("folder"):
        entry["folder"] = prev["folder"]     # keep the user's folder across re-index/refresh
    docs = [d for d in load_uploads() if d.get("file") != entry["file"]]
    docs.insert(0, entry)
    save_uploads(docs)

def remove_upload_record(file_name: str):
    save_uploads([d for d in load_uploads() if d.get("file") != file_name])

# ── User-created folders (organize uploaded knowledge sources) ───────────────────
FOLDERS_FILE = STORAGE_DIR / "folders.json"

def load_folders() -> list:
    if FOLDERS_FILE.exists():
        try:    return json.loads(FOLDERS_FILE.read_text())
        except: pass
    return []

def save_folders(names: list):
    FOLDERS_FILE.write_text(json.dumps(sorted({n.strip() for n in names if n and n.strip()}), indent=2))

# ── Background ingestion ───────────────────────────────────────────────────────
def ingest_file_bg(content: bytes, file_name: str, source_name: str, size_mb: float):
    PROCESSING_STATUS[file_name] = {"status": "processing", "chunks": 0}
    try:
        text        = extract_text(content, file_name)
        # Use smaller chunks for uploads (focused embeddings → better recall)
        text_chunks = chunk_text(text, size=UPLOAD_CHUNK_SIZE, overlap=UPLOAD_CHUNK_OVERLAP)
        if not text_chunks:
            raise ValueError("No extractable text found in this file.")

        # Grab the singleton ONCE at the start of ingestion.
        # All batch upserts go into the same object, so the query endpoint
        # (which also holds a reference to the same singleton) sees the
        # new chunks immediately after each upsert — no cache flush needed.
        col    = get_collection()
        stored = 0
        batch_ids, batch_docs, batch_metas, batch_embs = [], [], [], []

        for i, chunk in enumerate(text_chunks):
            try:
                e = embed(chunk)   # Ollama network call — outside the lock
            except Exception:
                continue           # skip un-embeddable chunks silently

            batch_ids.append(make_chunk_id(file_name, i))
            batch_docs.append(chunk)
            batch_metas.append({"source": source_name, "file": file_name,
                                 "chunk": i, "source_url": "",
                                 "origin": file_name, "origin_type": "file"})
            batch_embs.append(e)

            if len(batch_ids) >= 50:
                col.upsert(ids=batch_ids, documents=batch_docs,
                            metadatas=batch_metas, embeddings=batch_embs)
                stored += len(batch_ids)
                batch_ids.clear(); batch_docs.clear()
                batch_metas.clear(); batch_embs.clear()

        if batch_ids:
            col.upsert(ids=batch_ids, documents=batch_docs,
                        metadatas=batch_metas, embeddings=batch_embs)
            stored += len(batch_ids)

        # Do NOT reset the collection singleton — the same in-memory object
        # that just received the upserts IS what future queries will use.
        # Resetting would create a new client from SQLite which may not
        # have the writes flushed yet, making the new doc invisible to queries.

        if stored == 0:
            raise ValueError("Embedding failed for all chunks — is Ollama running?")

        PROCESSING_STATUS[file_name] = {"status": "indexed", "chunks": stored}
        cache_clear()                       # new content can change the best answer
        upsert_upload({"file": file_name, "name": source_name, "type": "file", "origin": file_name,
                       "size": f"{size_mb} MB", "status": "indexed",
                       "chunks": stored, "uploaded": time.strftime("%b %d, %Y")})

    except Exception as exc:
        err = str(exc)
        PROCESSING_STATUS[file_name] = {"status": "failed", "chunks": 0, "error": err}
        upsert_upload({"file": file_name, "name": source_name, "type": "file", "origin": file_name,
                       "size": f"{size_mb} MB", "status": "failed",
                       "chunks": 0, "uploaded": time.strftime("%b %d, %Y"),
                       "error": err})

# ── Knowledge from a Git repo or a web link (reuses ingest.py parsers) ───────
def _store_chunks(items: list, origin: str, origin_type: str) -> int:
    """Embed items [{source,file,chunk,document,source_url,spec_id?}] in batches and
    add them to the collection tagged with a shared `origin` so the whole source
    can be listed/deleted as a unit. Returns the number of chunks stored."""
    col, stored = get_collection(), 0
    for i in range(0, len(items), 64):
        batch = items[i:i + 64]
        embs  = ingest.embed_batch([it["document"] for it in batch])
        ids, docs, metas, es = [], [], [], []
        for it, e in zip(batch, embs):
            if e is None:
                continue
            ids.append("src_" + hashlib.md5(f"{origin}|{it['file']}|{it['chunk']}".encode()).hexdigest()[:16])
            docs.append(it["document"])
            metas.append({"source": it["source"], "file": it["file"], "chunk": it["chunk"],
                          "spec_id": it.get("spec_id", ""), "source_url": it.get("source_url", ""),
                          "origin": origin, "origin_type": origin_type})
            es.append(e)
        if ids:
            col.upsert(ids=ids, documents=docs, metadatas=metas, embeddings=es)
            stored += len(ids)
            PROCESSING_STATUS[origin] = {"status": "processing", "chunks": stored}
    return stored

def _html_to_text(html: str) -> str:
    import html as _h
    html = re.sub(r"(?is)<(script|style|noscript|head|nav|footer).*?</\1>", " ", html)
    html = re.sub(r"(?is)<br\s*/?>", "\n", html)
    html = re.sub(r"(?is)</(p|div|li|h[1-6]|tr|section|article)>", "\n", html)
    text = _h.unescape(re.sub(r"(?s)<[^>]+>", " ", html))
    return re.sub(r"\n{3,}", "\n\n", re.sub(r"[ \t]{2,}", " ", text)).strip()

def _fail(origin, label, typ, url, err):
    PROCESSING_STATUS[origin] = {"status": "failed", "chunks": 0, "error": err[:300]}
    upsert_upload({"file": origin, "name": label, "type": typ, "url": url, "origin": origin,
                   "size": typ, "status": "failed", "chunks": 0,
                   "uploaded": time.strftime("%b %d, %Y"), "error": err[:300]})

def _done(origin, label, typ, url, stored):
    PROCESSING_STATUS[origin] = {"status": "indexed", "chunks": stored}
    cache_clear()
    upsert_upload({"file": origin, "name": label, "type": typ, "url": url, "origin": origin,
                   "size": typ, "status": "indexed", "chunks": stored,
                   "uploaded": time.strftime("%b %d, %Y")})

def ingest_web_bg(url: str, origin: str, label: str):
    PROCESSING_STATUS[origin] = {"status": "processing", "chunks": 0}
    try:
        r = requests.get(url, timeout=40, headers={"User-Agent": "SynaptDI/1.0"})
        r.raise_for_status()
        if "pdf" in r.headers.get("content-type", "").lower() or url.lower().endswith(".pdf"):
            import fitz
            text = "\n".join(p.get_text() for p in fitz.open(stream=r.content, filetype="pdf"))
        else:
            text = _html_to_text(r.text)
        items = [{"source": label, "file": origin, "chunk": i, "document": p, "source_url": url}
                 for i, p in enumerate(chunk_text(text, size=900, overlap=150)) if len(p) >= 60]
        stored = _store_chunks(items, origin, "web")
        if stored == 0:
            raise ValueError("No readable text found at that URL.")
        _done(origin, label, "web", url, stored)
    except Exception as exc:
        _fail(origin, label, "web", url, str(exc))

def ingest_repo_bg(clone_url: str, origin: str, label: str):
    PROCESSING_STATUS[origin] = {"status": "processing", "chunks": 0}
    dest = UPLOADS_DIR / "_repos" / re.sub(r"[^A-Za-z0-9_.-]", "_", origin)
    try:
        shutil.rmtree(dest, ignore_errors=True)
        dest.parent.mkdir(parents=True, exist_ok=True)
        subprocess.run(["git", "clone", "--depth=1", clone_url, str(dest)],
                       check=True, capture_output=True, timeout=600)
        base   = re.sub(r"\.git$", "", clone_url.rstrip("/"))
        branch = "main"
        try:
            hb = subprocess.run(["git", "-C", str(dest), "rev-parse", "--abbrev-ref", "HEAD"],
                                capture_output=True, timeout=15)
            branch = hb.stdout.decode().strip() or "main"
        except Exception:
            pass
        blob = f"{base}/blob/{branch}" if base.startswith("http") else ""
        spec_files = []
        for ext in ("*.json", "*.yaml", "*.yml"):
            spec_files.extend(dest.rglob(ext))
        md_files = [f for f in dest.rglob("*.md") if not ingest.is_junk_markdown(f.name)]
        seen_hash, items = set(), []
        for f in spec_files + md_files:
            if any(p in f.parts for p in (".git", "node_modules", "__pycache__")):
                continue
            if "test" in f.name.lower() and f.suffix != ".md":
                continue
            parsed = ingest.parse_openapi_spec(f) if f.suffix != ".md" else ingest.parse_markdown(f)
            if not parsed or not parsed["text"].strip():
                continue
            h = hashlib.md5(parsed["text"].encode("utf-8", "ignore")).hexdigest()
            if h in seen_hash:
                continue
            seen_hash.add(h)
            src, sid = f"{parsed['title']} {parsed['version']}".strip(), ingest.extract_spec_id(f.name)
            surl = f"{blob}/{'/'.join(f.relative_to(dest).parts)}" if blob else ""
            for ci, ch in enumerate(ingest.chunk_text(parsed["text"])):
                if len(ch) >= 60:
                    items.append({"source": src, "file": f.name, "chunk": ci,
                                  "document": ch, "source_url": surl, "spec_id": sid})
        stored = _store_chunks(items, origin, "repo")
        if stored == 0:
            raise ValueError("No OpenAPI specs or docs found in that repository.")
        _done(origin, label, "repo", clone_url, stored)
    except subprocess.CalledProcessError as exc:
        _fail(origin, label, "repo", clone_url, exc.stderr.decode()[:200] if exc.stderr else "git clone failed")
    except Exception as exc:
        _fail(origin, label, "repo", clone_url, str(exc))
    finally:
        shutil.rmtree(dest, ignore_errors=True)   # discard the clone; embeddings are already stored

class RepoReq(BaseModel):
    url: str
    label: Optional[str] = None

class WebReq(BaseModel):
    url: str
    label: Optional[str] = None

@app.post("/sources/repo")
def add_repo_source(req: RepoReq, background_tasks: BackgroundTasks, _admin: dict = Depends(require_admin)):
    url = (req.url or "").strip()
    if not (url.lower().startswith(("http://", "https://")) or url.startswith("git@")):
        raise HTTPException(400, "Provide an https git URL, e.g. https://github.com/org/repo.")
    name   = re.sub(r"\.git$", "", url.rstrip("/")).split("/")[-1] or "repo"
    origin = "repo:" + name
    label  = (req.label or name).strip()
    upsert_upload({"file": origin, "name": label, "type": "repo", "url": url, "origin": origin,
                   "size": "repo", "status": "processing", "chunks": 0,
                   "uploaded": time.strftime("%b %d, %Y")})
    PROCESSING_STATUS[origin] = {"status": "processing", "chunks": 0}
    background_tasks.add_task(ingest_repo_bg, url, origin, label)
    return {"origin": origin, "status": "processing"}

@app.post("/sources/weblink")
def add_web_source(req: WebReq, background_tasks: BackgroundTasks, _admin: dict = Depends(require_admin)):
    url = (req.url or "").strip()
    if not url.lower().startswith(("http://", "https://")):
        raise HTTPException(400, "Provide a valid http(s) URL.")
    label  = (req.label or url.split("//", 1)[-1][:60]).strip()
    origin = "web:" + hashlib.md5(url.encode()).hexdigest()[:12]
    upsert_upload({"file": origin, "name": label, "type": "web", "url": url, "origin": origin,
                   "size": "web", "status": "processing", "chunks": 0,
                   "uploaded": time.strftime("%b %d, %Y")})
    PROCESSING_STATUS[origin] = {"status": "processing", "chunks": 0}
    background_tasks.add_task(ingest_web_bg, url, origin, label)
    return {"origin": origin, "status": "processing"}

# ── RAG ────────────────────────────────────────────────────────────────────────
# Two prompts, chosen by retrieval confidence.
#
# When the user names a spec we actually have indexed, the answer IS in the
# context — so we DON'T offer a refusal option. Small models (llama3.2 3B)
# over-refuse: a refusal sentence makes them bail even on clearly-relevant
# context (verified directly). For these high-confidence queries we force a
# grounded answer instead.
#
# For open-ended / non-spec queries we keep an out-of-domain escape hatch,
# because nomic-embed-text distances can't reliably separate in-domain from
# out-of-domain (an off-topic question can still land at a low distance).
_PROMPT_INSTRUCTIONS = """You answer the QUESTION using the CONTEXT below. Each source is labelled with its number and title, e.g. "[Source 2: TMF620 — Product Catalog]".

Write like a sharp, friendly expert explaining something to a teammate — clear, warm, and genuinely insightful, never robotic or one-word terse. Open with the direct answer in the very first sentence (no preamble, no label, no "Answer:"/"Summary:" heading). Then teach: explain what the key fields or concepts are for, how the pieces fit together, and any gotcha worth knowing — the context an expert would add so the reader actually understands, not just a bare list. A good answer usually has a one-line direct opener, a short plain-English explanation, and a structured detail block. Make every sentence count (no filler, no restating the question, no "In summary"/"I hope this helps" endings), but never be so brief that the answer carries no insight. Let the depth match the question.

Shape the answer to the question:
- Fields, attributes, query parameters, events, or status codes → a Markdown table such as `| Field | Required | Description |`, marking each mandatory or optional, and give every row a real plain-English description of what it is for (not just its name).
- Comparing two or more things (e.g. specs) → a Markdown table with one column each, then a one-line bottom line on which to use when.
- "How do I…" or how to call an operation → numbered steps plus a fenced code block with a concrete, runnable example (curl by default, or the language asked for) built from the real path and fields in the context.
- Anything else → short, scannable bullets or brief paragraphs.

Every specific field name, path, status code, and value must come from the CONTEXT — never invent them — but do explain and connect them in your own words so the answer teaches rather than just lists. Cite inline with the bracketed source number: put [n] right after the fact it backs (combine like [2][3]); cite only sources you used. Synthesize across related chunks. Answer confidently ONLY when the CONTEXT clearly supports the specific API, spec, version, field, operation, or capability named in the QUESTION. If the CONTEXT does not contain evidence for the exact thing asked about, say you don't have authoritative evidence for it rather than inferring it from a different, neighbouring, or unrelated API — never fabricate fields, operations, versions, or component relationships to fill a gap. Never open with filler like "Based on the provided context" or "According to the sources".

When asked what fields a resource requires, this almost always means what you must SUBMIT to create one: read the required list from the resource's _Create schema (e.g. ProductOrder_Create → productOrderItem), not the response resource's server-assigned id, not a *Ref or *_Update schema, and not an unrelated EventSubscription/Hub/Event/Error schema. Name the parent object a required field sits on (e.g. each productOrderItem needs an action and a productOffering) when the context shows it, and never answer with a single bare field name when you can explain what it is and where it lives. If the QUESTION names a specific version, answer for THAT major version only and do not substitute another major version's fields. If the CONTEXT mixes versions, state which version your answer reflects; never silently pick "the latest"."""

_GUARD_LINE = """

If the CONTEXT does not address the question, or the question is outside the telecom/standards domain entirely (general knowledge, geography, trivia, current events), do NOT answer from your own knowledge — reply with exactly:
"I don't have enough information in the current knowledge base to answer this. The relevant document may not be indexed yet." """

def _system_prompt(guarded: bool) -> str:
    """Build the system prompt, injecting the admin-configured domain persona."""
    persona = (load_chat_config().get("persona") or DEFAULT_CHAT_CONFIG["persona"]).strip()
    head = (f"You are SynaptDI (Synapt Domain Intelligence), an expert and friendly assistant for {persona}. "
            "You explain things clearly and helpfully, like a knowledgeable colleague who wants the reader to truly get it.")
    return f"{head}\n\n{_PROMPT_INSTRUCTIONS}" + (_GUARD_LINE if guarded else "")

# ── General-RAG answer shaping ────────────────────────────────────────────────────────
# The generic "Shape the answer" rules in _PROMPT_INSTRUCTIONS key off what's in the
# retrieved CONTEXT ("Fields, attributes → a table"), so an IDENTITY question ("what is
# TMF620?") whose chunks happen to be schema-heavy gets answered as a field/operations
# audit instead of a definition. This classifies the QUESTION's intent up front and, for
# definitional questions, injects an instruction that constrains answer SCOPE — the
# retrieved evidence is still the only source of facts; this is shaping, not invention.
# SCHEMA_GOVERNANCE never reaches here — it's owned by the deterministic route before RAG.
_SHAPE_IDENTITY = re.compile(
    r"^\s*(what\s+is|what'?s|what\s+are|whats)\b|^\s*(explain|describe|define)\b|"
    r"\btell\s+me\s+about\b|\b(overview|summary|introduction)\s+of\b|"
    r"^\s*what\s+does\b.*\bdo\b|\bwho\s+is\b|\bwhat\s+kind\s+of\b",
    re.I,
)
_SHAPE_COMPARISON = re.compile(
    r"\b(difference|differences|compare|comparison|versus|vs\.?)\b|\b\w+\s+vs\.?\s+\w+\b", re.I,
)
_SHAPE_HOWTO = re.compile(
    r"^\s*how\s+(do|to|can|would|should|does)\b|\bhow\s+do\s+i\b|\bexample\s+of\b|"
    r"\bshow\s+me\s+how\b|\bhow\s+would\s+i\b|\bwrite\s+(a|an|the)\b.*\b(request|call|query)\b",
    re.I,
)


def _answer_shape(question: str) -> str:
    """Lightweight general-intent classifier for the RAG answer's SHAPE (never routing):
    IDENTITY, COMPARISON, HOWTO, or GENERAL. Comparison is checked before identity so
    "what is the difference between X and Y" is COMPARISON, and how-to before identity so
    "how do I …" wins over an "explain"-like opener. No API-id or domain hardcoding."""
    q = question or ""
    if _SHAPE_HOWTO.search(q):
        return "HOWTO"
    if _SHAPE_COMPARISON.search(q):
        return "COMPARISON"
    if _SHAPE_IDENTITY.search(q):
        return "IDENTITY"
    return "GENERAL"


def _answer_shape_directive(question: str) -> str:
    """Per-question shaping instruction appended to the RAG system prompt. Empty for
    GENERAL (keep existing behaviour). Overrides the generic context-driven shaping rules
    for definitional/comparison/how-to intents."""
    shape = _answer_shape(question)
    if shape == "IDENTITY":
        return ("\n\nANSWER SHAPE — IDENTITY / DEFINITION (overrides the generic shaping rules "
                "above): the user asked WHAT this is, not for its schema. Lead with a one-sentence "
                "definition of what it is and its primary purpose/domain, then a short plain-English "
                "explanation of its role and, only if grounded in the CONTEXT, 3–5 core capabilities "
                "as brief bullets. You MAY note the current/latest version if the context shows it. "
                "Do NOT emit a Field/Required/Description (or any property) table, do NOT enumerate "
                "schema attributes, do NOT include curl/request/code examples, and do NOT turn this "
                "into a schema audit or a version-by-version comparison — unless the user explicitly "
                "asked for fields, an example, or a comparison. Keep it a crisp definitional overview.")
    if shape == "COMPARISON":
        return ("\n\nANSWER SHAPE — COMPARISON: contrast the named items directly. Represent EVERY "
                "named item using its own grounded evidence; if the context lacks evidence for one "
                "side, say so explicitly rather than inventing it. A compact comparison table or "
                "parallel bullets, then a one-line bottom line on when to use which. Do not pad with "
                "a full field/schema audit of either side unless asked.")
    if shape == "HOWTO":
        return ("\n\nANSWER SHAPE — HOW-TO / USAGE: give the concrete steps, and when the context "
                "supports it a short runnable example (endpoint path, query parameters, or a curl "
                "snippet) is appropriate and welcome. Keep it focused on accomplishing the task.")
    return ""

def _source_label(meta: dict) -> str:
    """Prefix the spec number so the model connects e.g. 'TMF622' (in the
    question) with 'Product Ordering' (the spec's title) — without it the model
    says the named spec 'is not in the context' and refuses."""
    sid, src = meta.get("spec_id"), meta.get("source", "Unknown")
    return f"{sid} — {src}" if sid else src

def _format_history(history) -> str:
    if not history:
        return ""
    turns = []
    for m in history[-6:]:
        c = (m.get("content") or "").strip()
        if not c:
            continue
        who = "User" if m.get("role") == "user" else "Assistant"
        turns.append(f"{who}: {c[:600]}")
    if not turns:
        return ""
    return ('RECENT CONVERSATION (use only to resolve references like "it"/"that"; '
            "the answer itself must still come from CONTEXT):\n" + "\n".join(turns) + "\n\n")

def _retrieval_query(question: str, history) -> str:
    """Fold the previous user turn into the retrieval query so follow-ups like
    'what about its events?' still retrieve the right spec."""
    if history:
        for m in reversed(history):
            if m.get("role") == "user" and (m.get("content") or "").strip():
                return f"{m['content'].strip()}\n{question}".strip()
    return question

def _ordered_sources(chunks: list) -> tuple:
    """Unique source names in first-seen order, a {name: 1-based number} map, and a
    per-chunk number list. Shared by build_prompt and build_sources so the inline
    [n] citations the model writes line up with the sources list sent to the client."""
    order, index, nums = [], {}, []
    for c in chunks:
        src = c["metadata"].get("source", "Unknown")
        if src not in index:
            order.append(src)
            index[src] = len(order)
        nums.append(index[src])
    return order, index, nums

def _unique_specs(chunks: list) -> list:
    """Ordered unique spec ids present in the retrieved chunks (for the retrieval trace)."""
    seen = []
    for c in chunks:
        sid = c["metadata"].get("spec_id")
        if sid and sid not in seen:
            seen.append(sid)
    return seen

# ── Spec version diff ───────────────────────────────────────────────────────────
_DIFF_INTENT = re.compile(r"\b(difference|differences|diff|changed|changes|compare|comparison|vs|versus|migrat\w*)\b", re.I)
_VER_TOKEN   = re.compile(r"\bv?(\d+)(?:\.\d+)*\b")

def _file_version(fname: str) -> str:
    """Major version from a spec filename, e.g. 'TMF622-...-v4.0.0.swagger.json' → '4'."""
    m = re.search(r"v(\d+)(?:\.\d+)*", fname or "", re.I)
    return m.group(1) if m else ""

def version_diff_plan(question: str, q_emb: list):
    """When the question asks what changed between versions of ONE spec, retrieve
    context from each version separately and return (chunks, extra_instruction);
    None means 'not a version-diff question — use the normal path'."""
    if not _DIFF_INTENT.search(question or ""):
        return None
    sids = detect_spec_ids(question)
    if len(sids) != 1:
        return None                                  # cross-spec comparisons use the normal path
    col = get_collection()
    files = sorted(get_spec_map(col).get(sids[0]) or [])
    by_ver: dict = {}
    for f in files:
        v = _file_version(f)
        if v:
            by_ver.setdefault(v, []).append(f)
    if len(by_ver) < 2:
        return None                                  # only one version indexed
    asked, seen = [], set()
    for v in _VER_TOKEN.findall(question):
        if v in by_ver and v not in seen:
            asked.append(v); seen.add(v)
    vers = asked[:2] if len(asked) >= 2 else sorted(by_ver, key=int)[-2:]
    chunks = []
    for v in vers:
        for h in col.query(q_emb, 6, where_in=("file", by_ver[v]))[:4]:
            chunks.append({"document": h["document"], "metadata": h["metadata"],
                           "distance": h["distance"], "upload": False, "eff": h["distance"]})
    if not chunks:
        return None
    lo, hi = sorted(vers, key=int)
    extra = (f"\n\nThis is a VERSION COMPARISON question about {sids[0]}. The context holds chunks from "
             f"v{lo} and v{hi} of the same spec — each source label names its version. Produce a structured "
             f"comparison: a Markdown table | Aspect | v{lo} | v{hi} | covering the fields, schemas, operations "
             "and behaviors that actually differ in the context, citing the version-specific source for each row, "
             "then one bottom line on what a team migrating between them should watch. State only differences the "
             "context shows — never invent a change.")
    return chunks, extra

def _project_extra(instructions: str) -> str:
    """Standing project notes/memory → an instruction block prepended to the system prompt
    so a chat inside a Project carries that project's context. Facts still come from CONTEXT."""
    t = (instructions or "").strip()
    if not t:
        return ""
    return ("\n\nPROJECT CONTEXT — this chat belongs to a user-created project with these standing "
            "notes/memory. Keep them in mind and stay consistent with them, but still ground every "
            "factual claim in the CONTEXT sources below:\n" + t[:2000])

def build_prompt(question: str, chunks: list[dict], guarded: bool, history=None, extra: str = "") -> str:
    sysp = _system_prompt(guarded) + extra
    _, index, _ = _ordered_sources(chunks)
    # Number by SOURCE (not by chunk) so multiple chunks of the same spec share one
    # citation number, and so [n] resolves to sources[n-1] on the client.
    ctx  = "\n\n".join(
        f"[Source {index[c['metadata'].get('source', 'Unknown')]}: {_source_label(c['metadata'])}]\n{c['document']}"
        for c in chunks
    )
    return f"{sysp}\n\n{_format_history(history)}CONTEXT:\n{ctx}\n\nQUESTION: {question}\n\nANSWER:"

def _gen_options(num_predict: int) -> dict:
    """Shared Ollama generation options. temperature 0 for determinism; full context
    so quality is never reduced; optional CPU thread pinning."""
    opts = {"temperature": 0.0, "num_predict": num_predict, "num_ctx": NUM_CTX,
            "stop": ["QUESTION:", "CONTEXT:"]}
    if NUM_THREAD > 0:
        opts["num_thread"] = NUM_THREAD
    return opts

def generate_answer(prompt: str, model: str = "", num_predict: int = DEEP_NUM_PREDICT) -> str:
    with _llm_sem:
        r = requests.post(f"{OLLAMA_URL}/api/generate",
                          json={"model": model or LLM_MODEL, "prompt": prompt, "stream": False,
                                "keep_alive": KEEP_ALIVE, "options": _gen_options(num_predict)},
                          timeout=600)
        r.raise_for_status()
        return r.json()["response"].strip()

# ── Grounded generation bridge for deterministic spec facts ────────────────────────
# A required-fields question is resolved deterministically first (spec_facts), then
# handed to the LLM as locked ground truth for a rich explanation, then integrity-
# validated before it's ever returned: DETERMINISTIC RESOLUTION -> VERIFIED FACT
# OBJECT -> GROUNDED LLM GENERATION -> INTEGRITY VALIDATION -> RESPONSE. No RAG
# chunks are mixed in here — the fact object already carries the schema-sourced
# detail (field descriptions, nested requirements) the model needs, and skipping
# retrieval sidesteps any version-mixing risk from the general corpus (Phase 5).
def _build_grounded_prompt(question: str, extra: str, history=None) -> str:
    return f"{_system_prompt(False)}{extra}\n\n{_format_history(history)}QUESTION: {question}\n\nANSWER:"

def _grounded_required_fields_answer(fact: dict, req, model: str, npred: int) -> dict:
    def _gen(q, extra):
        return generate_answer(_build_grounded_prompt(q, extra, req.history), model, npred)
    return spec_facts.grounded_answer(req.question, fact, _gen)

def _grounded_comparison_answer(comparison: dict, req, model: str, npred: int) -> dict:
    def _gen(q, extra):
        return generate_answer(_build_grounded_prompt(q, extra, req.history), model, npred)
    return spec_facts.grounded_comparison_answer(req.question, comparison, _gen)

def _grounded_property_answer(fact: dict, req, model: str, npred: int) -> dict:
    def _gen(q, extra):
        return generate_answer(_build_grounded_prompt(q, extra, req.history), model, npred)
    return spec_facts.grounded_property_answer(req.question, fact, _gen)

def _grounded_aggregation_answer(fact: dict, req, model: str, npred: int) -> dict:
    def _gen(q, extra):
        return generate_answer(_build_grounded_prompt(q, extra, req.history), model, npred)
    return spec_facts.grounded_aggregation_answer(req.question, fact, _gen)

_HIGH_CONF = {"level": "high", "score": 100, "strong": 1}
_LOW_CONF  = {"level": "low", "score": 0, "strong": 0}

def _resolve_governance_route(req, model: str, npred: int):
    """Runs the full governance router — required-fields (single-version or
    comparison), property claims, and corpus-wide aggregation, resolved or explicitly
    unresolved/refused — and returns (result, confidence) ready for
    _fact_response_payload/_fact_stream_response, or None if the question isn't a
    governance-shaped question at all (caller falls through to ODA, then general RAG).
    This is the ONE place that dispatches on fact_type/status, called identically by
    both /query and /query/stream, so they can never diverge on routing, resolved
    version, exhaustiveness, or provenance (Phase 7: a governance-shaped question
    either resolves deterministically or comes back explicitly UNRESOLVED/NON_EXHAUSTIVE
    — `spec_facts.route_governance_question` never returns a silent miss for one)."""
    rf = spec_facts.route_governance_question(req.question)
    if rf is None:
        return None
    raw = spec_facts.wants_raw(req.question)
    ft = rf.get("fact_type")

    if ft == "REQUIRED_FIELDS_COMPARISON":
        if rf["status"] == "UNRESOLVED":
            return spec_facts.unresolved_comparison_answer(rf), _LOW_CONF
        result = (spec_facts.raw_comparison_answer(rf) if raw
                  else _grounded_comparison_answer(rf, req, model, npred))
        return result, _HIGH_CONF

    if ft == "REQUIRED_FIELDS":
        if rf["status"] == "UNRESOLVED":
            return spec_facts.unresolved_answer(rf), _LOW_CONF
        result = (spec_facts.raw_fact_answer(rf) if raw
                  else _grounded_required_fields_answer(rf, req, model, npred))
        return result, _HIGH_CONF

    if ft in ("PROPERTY_EXISTENCE", "PROPERTY_REQUIRED", "PROPERTY_OPTIONAL"):
        result = (spec_facts.raw_property_fact_answer(rf) if raw
                  else _grounded_property_answer(rf, req, model, npred))
        return result, _HIGH_CONF

    if ft in ("SCHEMA_AGGREGATION", "REQUIRED_PROPERTY_AGGREGATION"):
        if rf["status"] == "UNRESOLVED":
            return spec_facts.unresolved_aggregation_answer(rf), _LOW_CONF
        if not rf["exhaustive"]:
            # Phase 5: never let the LLM near a corpus-wide claim the system can't fully
            # verify — the deterministic refusal-with-partial-detail IS the answer.
            return spec_facts.render_aggregation_fact_markdown(rf), _LOW_CONF
        result = (spec_facts.raw_aggregation_fact_answer(rf) if raw
                  else _grounded_aggregation_answer(rf, req, model, npred))
        return result, _HIGH_CONF

    # GOVERNANCE_REFUSAL (or any other unresolved shape): unmistakably a TM Forum
    # schema-fact question, but nothing above resolved anything — never fall through
    # to general RAG for it (Phase 7).
    return rf, _LOW_CONF

def _fact_response_payload(req, scope, mode, start, who, result: dict, conf: dict) -> "QueryResponse":
    """Shared /query response for a deterministically-resolved fact (required-fields,
    comparison, or ODA) — same cache/log handling either way; confidence is the
    caller's call (high for a resolved fact, low for an explicit unresolved result)."""
    _log_query(req.question, True, who)
    if not req.history and not req.no_cache:
        cache_put(req.question, scope, mode, result["answer"], result["sources"], conf,
                  project_instructions=req.project_instructions)
    return QueryResponse(answer=result["answer"], sources=result["sources"],
                         latency_ms=int((time.time() - start) * 1000),
                         chunks_retrieved=1, confidence=conf)

def _fact_stream_response(req, result: dict, conf: dict, start, who) -> StreamingResponse:
    """Shared /query/stream response for a deterministically-resolved fact. The answer is
    already fully resolved (and, for required-fields, already integrity-validated) before
    this is called, so — same as the short-circuit this replaces — it streams as one
    chunk rather than token-by-token. Caching is the caller's job (it needs scope/mode)."""
    def gen():
        _log_query(req.question, True, who)
        yield json.dumps({"type": "context", "sources": result["sources"],
                          "specs": [result["tmf"]] if result.get("tmf") else [],
                          "confidence": conf}) + "\n"
        yield json.dumps({"type": "token", "text": result["answer"]}) + "\n"
        yield json.dumps({"type": "done", "sources": result["sources"], "confidence": conf,
                          "latency_ms": int((time.time() - start) * 1000),
                          "chunks_retrieved": 1}) + "\n"
    return StreamingResponse(gen(), media_type="application/x-ndjson")

# ── Schemas ────────────────────────────────────────────────────────────────────
class QueryRequest(BaseModel):
    question: str
    top_k:    int           = TOP_K
    scope:    str           = "all"   # "all" = KB + your docs · "kb" = TM Forum only · "docs" = your uploads only
    history:  Optional[list] = None   # [{role, content}] recent turns, for follow-up questions
    standards_filter: Optional[list] = None
    mode:     str           = "deep"  # "deep" = full 8B model · "fast" = small model for quick lookups
    no_cache: bool          = False   # Regenerate sets this → skip the cache read, force a fresh answer
    project_instructions: str = ""    # standing notes/memory of the project this chat belongs to

class Source(BaseModel):
    name: str; file: str; chunk: int; preview: str; url: str = ""; upload: bool = False

class QueryResponse(BaseModel):
    answer: str; sources: list; latency_ms: int; chunks_retrieved: int; cached: bool = False; confidence: dict = {}

class FollowupReq(BaseModel):
    question: str = ""
    answer:   str = ""

# ══════════════════════════════════════════════════════════════════════════════
# ENDPOINTS
# ══════════════════════════════════════════════════════════════════════════════

# ── Health / Stats ─────────────────────────────────────────────────────────────
@app.get("/health")
def health():
    try:
        col    = get_collection()
        r      = requests.get(f"{OLLAMA_URL}/api/tags", timeout=5)
        models = [m["name"] for m in r.json().get("models", [])]
        return {"status":"ok","chunks_indexed":col.count(),"ollama_models":models,
                "llm_model":LLM_MODEL,"embed_model":EMBED_MODEL,
                "collection":vectorstore.COLLECTION,
                "collection_source":("explicit" if vectorstore.COLLECTION_IS_EXPLICIT else "default")}
    except Exception as e:
        raise HTTPException(503, str(e))

@app.get("/stats")
def stats():
    col = get_collection()
    return {"chunks_indexed": col.count(), "collection": vectorstore.COLLECTION,
            "collection_source": ("explicit" if vectorstore.COLLECTION_IS_EXPLICIT else "default")}

@app.get("/coverage")
def coverage():
    """What the knowledge base actually covers — specs grouped by domain (for onboarding)."""
    col = get_collection()
    try:
        smap = get_spec_map(col)
    except Exception:
        smap = {}
    groups: dict = {}
    for sid, files in smap.items():
        f0 = sorted(files)[0] if files else ""
        dom = _kb_category(sid, f0, "")
        groups.setdefault(dom, set()).add(sid)
    domains = [{"domain": d, "specs": sorted(groups[d])} for d in sorted(groups)]
    return {"chunks": col.count(), "spec_count": len(smap), "domains": domains}

# ── Query ──────────────────────────────────────────────────────────────────────
# Schema-aware re-ranking. nomic-embed distances alone can rank an unrelated schema
# that literally says "required: [id, callback]" (e.g. EventSubscription) above the
# resource the user actually named (e.g. ProductOrder). We nudge the ranking by
# matching the schema name embedded in each chunk to the resource in the question.
_SCHEMA_RE    = re.compile(r"Schema:\s*([A-Za-z0-9_]+)")
_FIELD_INTENT = re.compile(r"\b(mandatory|required|require[sd]?|fields?|attributes?|propert(?:y|ies)|schema)\b", re.I)
_FOCUS_STOP   = {"what","whats","which","whose","does","do","did","the","a","an","of","for","in","on","is",
                 "are","to","how","require","requires","required","mandatory","optional","fields","field",
                 "attributes","attribute","property","properties","schema","schemas","list","show","tell",
                 "me","about","and","or","resource","entity","object","model","open","api","apis","need",
                 "needs","needed","have","has","with","please","give","when","creating","create"}
_GENERIC_SCHEMAS = ("eventsubscription","hub","event","notification","error","extensibleerror","meta",
                    "entityref","timeperiod","money","quantity","note","characteristic","attachmentref")
LEX_WEIGHT = 0.12   # hybrid re-rank: how much a full keyword match can pull a chunk up
DOMAIN_GUARD_DIST = 0.34   # best-hit distance worse than this + zero keyword overlap ⇒ off-domain

def _schema_name(doc: str) -> str:
    m = _SCHEMA_RE.search(doc or "")
    return m.group(1) if m else ""

def _resource_focus(question: str) -> dict:
    """Pull the resource the user is asking about out of the question, e.g.
    'mandatory fields of TMF622 Product Order' -> target 'productorder', plus the
    set of content keywords used for the lexical half of the hybrid re-rank."""
    cleaned = re.sub(r"\bTMF\s?\d{3,4}\b", " ", question or "", flags=re.I)
    words   = re.findall(r"[A-Za-z][A-Za-z0-9]*", cleaned)
    keep    = [w for w in words if w.lower() not in _FOCUS_STOP and len(w) > 1]
    terms   = sorted({w.lower() for w in words if len(w) > 2 and w.lower() not in _FOCUS_STOP})
    return {"target": "".join(keep).lower(),
            "is_field": bool(_FIELD_INTENT.search(question or "")),
            "terms": terms}

def _rank_adjust(focus: dict, doc: str, dist: float) -> float:
    """Hybrid re-rank → effective distance (lower = better). Blends the dense embedding
    distance with two sparse signals: (a) a schema-name match to the named resource
    (preferring its _Create variant for 'fields' questions, demoting *Ref/*_Update and
    unrelated generic schemas), and (b) lexical overlap with the question's keywords."""
    eff = dist
    target = focus.get("target") or ""
    sn = _schema_name(doc).lower().replace("_", "") if target else ""
    if target and sn:
        is_field = focus.get("is_field")
        if is_field and any(sn == g or sn.startswith(g) for g in _GENERIC_SCHEMAS):
            eff += 0.25                           # unrelated notification/error schema on a fields question
        elif is_field and target in sn and (sn.endswith("ref") or sn.endswith("update")):
            eff += 0.15                           # *Ref / *_Update aren't where required-to-create lives
        elif sn == target + "create":
            eff -= 0.22                           # the _Create / FVO variant — exactly what you must submit
        elif sn == target:
            eff -= 0.16                           # the main resource schema
        elif target in sn or sn in target:
            eff -= 0.10                           # related schema (e.g. ProductOrderItem)
    # Lexical (sparse) signal: reward chunks that literally contain the query's keywords.
    terms = focus.get("terms") or []
    if terms:
        dl = doc.lower()
        hits = sum(1 for t in terms if t in dl)
        if hits:
            eff -= LEX_WEIGHT * (hits / len(terms))
    return eff

def retrieve(question: str, q_emb: list, top_k: int, scope: str, hints: dict = None):
    """Shared retrieval for /query and /query/stream.
    Returns (chunks, confident, empty_msg). chunks == [] means nothing relevant;
    empty_msg holds the user-facing fallback text in that case.
    `hints` (Phase 7B, from knowledge_router) may carry: spec_ids to seed, a requested
    major version to ISOLATE (excluding other majors), and canonical_only to keep a
    named-spec question off ODA-Canvas operator prose."""
    col = get_collection()
    total = col.count()
    if total == 0:
        return [], False, NO_INFO_MSG

    hints          = hints or {}
    hint_specs     = [s.upper() for s in (hints.get("spec_ids") or [])]
    hint_major     = hints.get("major")
    canonical_only = bool(hints.get("canonical_only"))
    added     = [d.get("origin") or d.get("file") for d in load_uploads() if d.get("status") == "indexed"]
    added_set = set(added)
    focus     = _resource_focus(question)
    seen:    set  = set()
    chunks:  list = []
    matched: list = []

    def add(d, m, dist, upload=False, eff=None) -> bool:
        key = (m.get("file", ""), m.get("chunk", 0))
        if key in seen:
            return False
        seen.add(key)
        chunks.append({"document": d, "metadata": m, "distance": dist, "upload": upload,
                       "eff": dist if eff is None else eff})
        return True

    if scope == "docs":
        # Answer ONLY from the user's uploaded documents.
        if added:
            for h in col.query(q_emb, min(top_k * 2, total), where_in=("origin", added)):
                if len(chunks) >= top_k:
                    break
                if h["distance"] < DOC_SCOPE_MAX:
                    add(h["document"], h["metadata"], h["distance"], upload=True)
    else:
        # 1. Spec-targeted retrieval — guarantee named specs are present.
        spec_map = get_spec_map(col)
        spec_ids = detect_spec_ids(question)
        for s in hint_specs:                          # Phase 7B: router-seeded specs (bare entities, links)
            if s not in spec_ids:
                spec_ids.append(s)
        matched  = [s for s in spec_ids if spec_map.get(s)]
        if matched:
            per_spec = max(top_k // len(matched), 2)
            for sid in matched:
                sfiles = sorted(spec_map[sid])
                if hint_major is not None:            # Phase 7B/7C: version isolation — drop other majors
                    sfiles = [f for f in sfiles if knowledge_link.major_of(f) == str(hint_major)]
                    if not sfiles:
                        continue                       # requested major absent → no evidence (never fall back)
                hits = col.query(q_emb, min(per_spec * 4, total), where_in=("file", sfiles))
                cands = [(_rank_adjust(focus, h["document"], h["distance"]), h["document"], h["metadata"], h["distance"])
                         for h in hits if h["distance"] < KB_DIST_MAX]
                cands.sort(key=lambda x: x[0])           # schema-aware re-rank within the named spec
                n = per_spec
                for eff, d, m, dist in cands:
                    if n <= 0:
                        break
                    if add(d, m, dist, eff=eff):
                        n -= 1

        # 2. Merge uploaded docs (boosted) + general KB by effective distance.
        pool: list = []
        if scope == "all" and added:
            for h in col.query(q_emb, min(top_k * 2, total), where_in=("origin", added)):
                if h["distance"] < UPLOAD_MAX_DIST:
                    pool.append({"d": h["document"], "m": h["metadata"], "dist": h["distance"], "upload": True,
                                 "eff": max(0.0, _rank_adjust(focus, h["document"], h["distance"]) - UPLOAD_BOOST)})
        for h in col.query(q_emb, min(top_k * 3, total)):
            m = h["metadata"]
            if m.get("origin", "") in added_set:        # user-added sources handled separately above
                continue
            cc = knowledge_link.content_class(m.get("file", ""))
            if cc in ("non_tmf_external", "test_collection"):
                continue                                # Phase 7B: MEF/postman never rank as canonical evidence
            if canonical_only and cc == "oda_canvas":
                continue                                # a named-spec question stays off Canvas operator prose
            if h["distance"] < KB_DIST_MAX:
                pool.append({"d": h["document"], "m": m, "dist": h["distance"], "upload": False,
                             "eff": _rank_adjust(focus, h["document"], h["distance"])})
        pool.sort(key=lambda c: c["eff"])
        for c in pool:
            if len(chunks) >= top_k:
                break
            add(c["d"], c["m"], c["dist"], upload=c["upload"], eff=c["eff"])

    if not chunks:
        if scope == "docs":
            msg = ("You haven't uploaded any documents yet — add one on the Documents page."
                   if not upload_files else
                   "I couldn't find anything relevant in your uploaded documents for that question.")
        else:
            msg = NO_INFO_MSG
        return [], False, msg

    # Domain guard: when retrieval is semantically weak AND the question shares zero
    # vocabulary with what it retrieved, it's off-domain trivia (e.g. "capital of
    # France" lands at 0.36 — closer than some real eTOM questions, so a distance
    # floor alone can't work). Feeding the LLM junk context makes it answer from
    # world knowledge; refuse instead. Measured: in-domain questions always share
    # at least one content term with their retrieved chunks.
    if scope != "docs" and not matched:
        best_raw = min(c["distance"] for c in chunks)
        terms = focus.get("terms") or []
        if terms and best_raw > DOMAIN_GUARD_DIST:
            # Count overlap only in semantically close chunks (raw ≤ 0.40) — the
            # lexical re-rank can boost a far chunk that shares one incidental word
            # ("capital" in a finance doc), and that must not defeat the guard.
            blob = " ".join(c["document"].lower() for c in chunks if c["distance"] <= 0.40)
            if not any(t in blob for t in terms):
                return [], False, NO_INFO_MSG

    chunks.sort(key=lambda c: c.get("eff", c["distance"]))   # schema-aware order → best match becomes [Source 1]
    best = min(c["distance"] for c in chunks)
    if scope == "docs":
        confident = best < UPLOAD_MAX_DIST
    else:
        confident = bool(matched) \
            or any(c["upload"] and c["distance"] < UPLOAD_CONFIDENT for c in chunks) \
            or best < CONFIDENT_DIST
    return chunks, confident, ""

def build_sources(chunks: list) -> list:
    """De-duplicated source list (JSON-serialisable), ordered to match the [Source n]
    numbering used in the prompt so inline [n] citations resolve to the right entry."""
    order, _, _ = _ordered_sources(chunks)
    first = {}
    for c in chunks:
        src = c["metadata"].get("source", "Unknown")
        if src not in first:
            first[src] = c
    sources = []
    rec = []
    for src in order:
        c = first[src]
        QUERY_HITS[src] = QUERY_HITS.get(src, 0) + 1
        rec.append((src, c.get("distance")))
        meta  = c["metadata"]
        otype = meta.get("origin_type", "")          # "file"/"repo"/"web" for user sources, "" for bundled KB
        sources.append({
            "name":    src,
            "file":    meta.get("file", ""),
            "chunk":   meta.get("chunk", 0),
            "preview": c["document"][:400] + ("..." if len(c["document"]) > 400 else ""),
            "url":     meta.get("source_url", ""),
            "upload":  bool(c.get("upload")),
            "origin_type": otype,
            "domain":  "" if otype else _kb_category(meta.get("spec_id", ""), meta.get("file", ""), ""),
        })
    _record_doc_hits(rec)
    return sources

def stream_ollama(prompt: str, model: str = "", num_predict: int = DEEP_NUM_PREDICT):
    """Yield answer text fragments from Ollama as they are generated.
    Generations are serialized via _llm_sem so concurrent users queue cleanly."""
    _llm_sem.acquire()
    try:
        with requests.post(f"{OLLAMA_URL}/api/generate",
                           json={"model": model or LLM_MODEL, "prompt": prompt, "stream": True,
                                 "keep_alive": KEEP_ALIVE, "options": _gen_options(num_predict)},
                           stream=True, timeout=600) as r:
            r.raise_for_status()
            for line in r.iter_lines():
                if not line:
                    continue
                try:    obj = json.loads(line)
                except Exception: continue
                frag = obj.get("response", "")
                if frag:
                    yield frag
                if obj.get("done"):
                    break
    finally:
        _llm_sem.release()

# ── Phase 7C: generic evidence-entailment guards (no hardcoded concepts) ──────────────
# A "which/what TMF API defines/represents X" question, when no exact entity was resolved,
# must NOT be answered from a mere lexical neighbour (e.g. "cryptocurrency wallets" matching
# a payment "DigitalWallet" doc). Require every distinctive concept term of the question to
# actually appear in the retrieved canonical evidence; if a concept is absent, abstain.
_DEFINE_INTENT = re.compile(
    r"\b(which|what)\b[^?]{0,60}\b(api|apis|specification|spec|standard)\b"
    r"|\b(define[sd]?|models?|represent[sd]?)\b", re.I)
_PREMISE_STOP = {
    "which", "what", "does", "tmf", "tmforum", "forum", "api", "apis", "spec", "specification",
    "standard", "define", "defines", "definition", "model", "models", "manage", "manages",
    "handle", "handles", "represent", "represents", "provide", "provides", "there", "this",
    "that", "exist", "exists", "have", "with", "from", "about", "used", "using", "should",
    "would", "could", "your", "mine", "real", "actual", "official", "form", "forums"}


def _premise_terms(question: str) -> list:
    return [t for t in re.findall(r"[a-z]{4,}", (question or "").lower()) if t not in _PREMISE_STOP]


def _false_premise_abstention(question: str, chunks: list, khints: Optional[dict]):
    """Return an abstention string when a define/relation question names a concept that is
    ABSENT from the retrieved evidence (only a lexical near-neighbour matched); else None."""
    if khints and khints.get("evidence_specs"):
        return None                                  # an exact entity was resolved — not this guard
    if not _DEFINE_INTENT.search(question or ""):
        return None
    terms = _premise_terms(question)
    if not terms:
        return None
    blob = " ".join((c.get("document") or "").lower() for c in chunks)
    missing = [t for t in terms
               if t not in blob and (t[:-1] if t.endswith("s") else t + "s") not in blob and t.rstrip("s") not in blob]
    if missing:
        return (f"I don't have authoritative TM Forum specification evidence that any API "
                f"represents **{' '.join(missing)}** — the retrieved material is only a lexical "
                f"neighbour, not a specification that defines that concept, so I won't affirm it.")
    return None


def _spec_identity_extra(khints: Optional[dict], chunks: list) -> str:
    """When a bare entity linked to an owning spec (e.g. 'fault'/'alarm' → TMF642) and that
    spec is actually present in the evidence, instruct the generator to state the TMF spec
    identity precisely. Generic — driven by the linked spec_id, not a hardcoded phrase."""
    if not khints:
        return ""
    specs = khints.get("evidence_specs") or khints.get("spec_ids")
    if not specs:
        return ""
    got = {(c["metadata"].get("spec_id") or "").upper() for c in chunks}
    named = [s for s in specs if s.upper() in got]
    if not named:
        return ""
    schema = khints.get("schema")
    obj = f"the '{schema}' resource" if schema else "this topic"
    return (f"\n\nThe question concerns {obj}, which the retrieved canonical evidence attributes to "
            f"the {named[0]} TM Forum specification. State that specification identity ({named[0]}) "
            f"explicitly and precisely in your answer; do not give a generic description.")


@app.post("/query", response_model=QueryResponse, tags=["Knowledge"])
def query(req: QueryRequest, authorization: Optional[str] = Header(None)):
    who = optional_user_email(authorization)
    if not req.question.strip():
        raise HTTPException(400, "Question cannot be empty")
    start = time.time()
    top_k = max(1, min(req.top_k, 20))
    scope = (req.scope or "all").lower()
    if scope not in ("all", "kb", "docs"):
        scope = "all"

    mode  = "fast" if (req.mode or "").lower() == "fast" else "deep"
    model = FAST_MODEL if mode == "fast" else LLM_MODEL
    npred = FAST_NUM_PREDICT if mode == "fast" else DEEP_NUM_PREDICT
    if mode == "fast":
        top_k = min(top_k, FAST_TOP_K)   # fewer chunks → smaller prompt → faster on CPU
    if not req.history and not req.no_cache:  # skip cache for follow-ups and Regenerate
        hit = cache_get(req.question, scope, mode, project_instructions=req.project_instructions)
        if hit:
            _log_query(req.question, True, who)
            return QueryResponse(answer=hit["answer"], sources=hit["sources"],
                                 latency_ms=int((time.time()-start)*1000),
                                 chunks_retrieved=0, cached=True, confidence=hit.get("confidence", {}))

    # Deterministic governance-fact resolution: structured questions ("what mandatory
    # fields does TMF622 Product Order require?", confirmations/challenges naming a
    # specific field — including one that doesn't exist, "which schemas contain/require
    # X" corpus-wide aggregation, and version comparisons) are resolved straight from
    # the canonical schema — correct by construction, never sampled from retrieval. An
    # explicit version that can't be found, or a corpus-wide claim whose enumeration
    # can't be completed exhaustively, resolves to an explicit unresolved/refused
    # result — never a silent fall-through to general RAG that then answers the
    # governance question generatively. Unless the user explicitly asked for raw,
    # unexplained output, a resolved fact is handed to the LLM as locked ground truth
    # for a rich explanation, integrity-validated before being returned; any violation
    # falls back to the exact deterministic answer. ODA component questions are
    # unaffected — still answered directly, as before.
    khints = None
    if scope in ("all", "kb"):
        # Knowledge Engine V2 (Phase 7B): deterministic ODA-contract routing + honest
        # abstention run FIRST — a named component/spec question must not be mis-grabbed by
        # the schema-required-fields router or the fuzzy ODA-mapping engine, and a named
        # entity must never be answered from a semantically-near unrelated spec.
        kdec = knowledge_router.route(req.question)
        khints = kdec.get("hints")     # capture early — used by the version guard below
        if kdec.get("kind") == "answer":
            return _fact_response_payload(req, scope, mode, start, who,
                                          {"answer": kdec["answer"], "sources": kdec["sources"]}, _HIGH_CONF)
        if kdec.get("kind") == "abstain":
            _log_query(req.question, False, who)
            return QueryResponse(answer=kdec["answer"], sources=[],
                                 latency_ms=int((time.time()-start)*1000), chunks_retrieved=0,
                                 confidence=_LOW_CONF)
        # Existing deterministic schema-required-fields engine handles "what mandatory fields
        # does TMF622 require?". Phase 7C version guard: if the user asked for an EXPLICIT major
        # but governance resolves a DIFFERENT major (it defaults to a single version), do NOT
        # return a wrong-version fact — fall through to version-isolated RAG for the asked major.
        routed = _resolve_governance_route(req, model, npred)
        if routed:
            req_major = str(khints["major"]) if (khints and khints.get("major") is not None) else None
            if req_major is not None:
                _rf = spec_facts.required_fields_fact(req.question)
                _rv = re.search(r"(\d+)", str((_rf or {}).get("version") or "")) if _rf else None
                if _rv and _rv.group(1) != req_major:
                    routed = None      # governance can't honour the requested version → RAG
            if routed:
                result, conf = routed
                return _fact_response_payload(req, scope, mode, start, who, result, conf)
        oda_fact = spec_facts.oda_answer(req.question)
        if oda_fact:
            return _fact_response_payload(req, scope, mode, start, who, oda_fact, _HIGH_CONF)

    rq = _retrieval_query(req.question, req.history)
    try:    q_emb = embed_query(rq)
    except Exception as e: raise HTTPException(503, f"Embedding failed: {e}")

    extra = ""
    plan = version_diff_plan(req.question, q_emb)
    if plan:
        chunks, extra = plan
        confident, empty = True, ""
    else:
        try:    chunks, confident, empty = retrieve(rq, q_emb, top_k, scope, hints=khints)
        except Exception as e: raise HTTPException(503, f"Retrieval failed: {e}")

    # Evidence validation (Phase 7B): when the router named specific specs the answer must
    # be grounded in, refuse rather than answer from unrelated retrieved content.
    if khints and khints.get("evidence_specs") and chunks:
        want = {s.upper() for s in khints["evidence_specs"]}
        got = {(c["metadata"].get("spec_id") or "").upper() for c in chunks}
        if not (want & got):
            _log_query(req.question, False, who)
            return QueryResponse(
                answer=("I found related material but no authoritative evidence for "
                        f"{', '.join(sorted(want))} specifically, so I won't answer from a "
                        "different specification. Try naming the exact API/version, or add the "
                        "spec on the Documents page."),
                sources=[], latency_ms=int((time.time()-start)*1000), chunks_retrieved=0,
                confidence=_LOW_CONF)

    if not chunks:
        _log_query(req.question, False, who)
        return QueryResponse(answer=empty, sources=[],
                             latency_ms=int((time.time()-start)*1000), chunks_retrieved=0)

    # Phase 7C: false-premise / lexical-neighbour guard (generic) — abstain before generating.
    fp = _false_premise_abstention(req.question, chunks, khints)
    if fp:
        _log_query(req.question, False, who)
        return QueryResponse(answer=fp, sources=[], latency_ms=int((time.time()-start)*1000),
                             chunks_retrieved=0, confidence=_LOW_CONF)

    ident = _spec_identity_extra(khints, chunks)   # Phase 7C: precise spec-identity grounding
    extra = _project_extra(req.project_instructions) + extra + ident + spec_facts.oda_grounding(req.question) + _answer_shape_directive(req.question)
    try:    answer = generate_answer(build_prompt(req.question, chunks, guarded=not confident, history=req.history, extra=extra), model, npred)
    except Exception as e: raise HTTPException(503, f"Generation failed: {e}")

    # If the LLM used its exact "no info" escape hatch, return no sources.
    if answer.lower().strip().startswith(NO_INFO_PREFIX):
        _log_query(req.question, False, who)
        return QueryResponse(answer=answer, sources=[],
                             latency_ms=int((time.time()-start)*1000), chunks_retrieved=0)

    _log_query(req.question, True, who)
    srcs = build_sources(chunks)
    conf = _confidence(chunks)
    if not req.history:
        cache_put(req.question, scope, mode, answer, srcs, conf, project_instructions=req.project_instructions)
    return QueryResponse(answer=answer, sources=srcs,
                         latency_ms=int((time.time()-start)*1000), chunks_retrieved=len(chunks), confidence=conf)


# ── Query (streaming, token-by-token) ────────────────────────────────────────
@app.post("/query/stream", tags=["Knowledge"])
def query_stream(req: QueryRequest, authorization: Optional[str] = Header(None)):
    who = optional_user_email(authorization)
    if not req.question.strip():
        raise HTTPException(400, "Question cannot be empty")
    start = time.time()
    top_k = max(1, min(req.top_k, 20))
    scope = (req.scope or "all").lower()
    if scope not in ("all", "kb", "docs"):
        scope = "all"
    mode  = "fast" if (req.mode or "").lower() == "fast" else "deep"
    model = FAST_MODEL if mode == "fast" else LLM_MODEL
    npred = FAST_NUM_PREDICT if mode == "fast" else DEEP_NUM_PREDICT
    if mode == "fast":
        top_k = min(top_k, FAST_TOP_K)   # fewer chunks → smaller prompt → faster on CPU
    if not req.history and not req.no_cache:
        hit = cache_get(req.question, scope, mode, project_instructions=req.project_instructions)
        if hit:
            def cached_gen():
                _log_query(req.question, True, who)
                yield json.dumps({"type": "token", "text": hit["answer"]}) + "\n"
                yield json.dumps({"type": "done", "sources": hit["sources"], "cached": True,
                                  "confidence": hit.get("confidence", {}),
                                  "latency_ms": int((time.time() - start) * 1000),
                                  "chunks_retrieved": 0}) + "\n"
            return StreamingResponse(cached_gen(), media_type="application/x-ndjson")

    # Deterministic governance-fact resolution (see /query for the full rationale) —
    # resolve (single-version, comparison, property claim, or corpus-wide aggregation;
    # resolved, explicitly unresolved, or explicitly refused) via the SAME router
    # /query uses, then either format raw or run the grounded-generation bridge, then
    # stream the final (already integrity-validated) answer as a single chunk.
    khints = None
    if scope in ("all", "kb"):
        # Knowledge Engine V2 (Phase 7C: streaming path mirrors /query so the chat UI gets
        # the same deterministic routing, abstention, version isolation and evidence gating).
        kdec = knowledge_router.route(req.question)
        khints = kdec.get("hints")
        if kdec.get("kind") == "answer":
            ans = {"answer": kdec["answer"], "sources": kdec["sources"]}
            if not req.history and not req.no_cache:
                cache_put(req.question, scope, mode, ans["answer"], ans["sources"], _HIGH_CONF,
                         project_instructions=req.project_instructions)
            return _fact_stream_response(req, ans, _HIGH_CONF, start, who)
        if kdec.get("kind") == "abstain":
            return _fact_stream_response(req, {"answer": kdec["answer"], "sources": []}, _LOW_CONF, start, who)
        routed = _resolve_governance_route(req, model, npred)
        if routed:
            req_major = str(khints["major"]) if (khints and khints.get("major") is not None) else None
            if req_major is not None:
                _rf = spec_facts.required_fields_fact(req.question)
                _rv = re.search(r"(\d+)", str((_rf or {}).get("version") or "")) if _rf else None
                if _rv and _rv.group(1) != req_major:
                    routed = None
            if routed:
                result, conf = routed
                if not req.history and not req.no_cache:
                    cache_put(req.question, scope, mode, result["answer"], result["sources"], conf,
                             project_instructions=req.project_instructions)
                return _fact_stream_response(req, result, conf, start, who)
        oda_fact = spec_facts.oda_answer(req.question)
        if oda_fact:
            if not req.history and not req.no_cache:
                cache_put(req.question, scope, mode, oda_fact["answer"], oda_fact["sources"], _HIGH_CONF,
                         project_instructions=req.project_instructions)
            return _fact_stream_response(req, oda_fact, _HIGH_CONF, start, who)

    rq = _retrieval_query(req.question, req.history)
    try:    q_emb = embed_query(rq)
    except Exception as e: raise HTTPException(503, f"Embedding failed: {e}")
    extra = ""
    plan = version_diff_plan(req.question, q_emb)
    if plan:
        chunks, extra = plan
        confident, empty = True, ""
    else:
        chunks, confident, empty = retrieve(rq, q_emb, top_k, scope, hints=khints)
    # Phase 7C: evidence + false-premise gates (mirror /query) before streaming generation.
    if khints and khints.get("evidence_specs") and chunks:
        if not ({s.upper() for s in khints["evidence_specs"]} & {(c["metadata"].get("spec_id") or "").upper() for c in chunks}):
            chunks = []
    if chunks:
        _fp = _false_premise_abstention(req.question, chunks, khints)
        if _fp:
            return _fact_stream_response(req, {"answer": _fp, "sources": []}, _LOW_CONF, start, who)
    _ident = _spec_identity_extra(khints, chunks)
    extra = _project_extra(req.project_instructions) + extra + _ident + spec_facts.oda_grounding(req.question) + _answer_shape_directive(req.question)

    def gen():
        if not chunks:
            _log_query(req.question, False, who)
            yield json.dumps({"type": "token", "text": empty}) + "\n"
            yield json.dumps({"type": "done", "sources": [],
                              "latency_ms": int((time.time() - start) * 1000),
                              "chunks_retrieved": 0}) + "\n"
            return
        conf = _confidence(chunks)
        # Tell the client what we're reading before the slow generation starts.
        yield json.dumps({"type": "context",
                          "sources": _ordered_sources(chunks)[0][:6],
                          "specs":   _unique_specs(chunks)[:6],
                          "confidence": conf}) + "\n"
        acc = ""
        try:
            for frag in stream_ollama(build_prompt(req.question, chunks, guarded=not confident, history=req.history, extra=extra), model, npred):
                acc += frag
                yield json.dumps({"type": "token", "text": frag}) + "\n"
        except Exception as e:
            yield json.dumps({"type": "error", "text": f"Generation failed: {e}"}) + "\n"
            return
        srcs = [] if acc.lower().strip().startswith(NO_INFO_PREFIX) else build_sources(chunks)
        _log_query(req.question, bool(srcs), who)
        if not req.history:
            cache_put(req.question, scope, mode, acc, srcs, conf, project_instructions=req.project_instructions)
        yield json.dumps({"type": "done", "sources": srcs,
                          "confidence": conf if srcs else {},
                          "latency_ms": int((time.time() - start) * 1000),
                          "chunks_retrieved": len(chunks)}) + "\n"

    return StreamingResponse(gen(), media_type="application/x-ndjson")

# ── Follow-up suggestions ────────────────────────────────────────────────────────
def generate_followups(question: str, answer: str) -> list:
    """Suggest 3 natural next questions from a Q&A pair. Best-effort; [] on any failure."""
    ans = (answer or "").strip()
    if not ans or ans.lower().startswith(NO_INFO_PREFIX):
        return []
    prompt = (
        "You help a user explore TM Forum / telecom standards. Given the question and answer "
        "below, suggest exactly 3 short, specific follow-up questions the user would naturally "
        "ask next. Each must be self-contained and under 12 words. Return ONLY the questions, "
        "one per line, with no numbering or bullets.\n\n"
        f"QUESTION: {question}\n\nANSWER: {ans[:1500]}\n\nFOLLOW-UPS:"
    )
    try:
        r = requests.post(f"{OLLAMA_URL}/api/generate",
                          json={"model": LLM_MODEL, "prompt": prompt, "stream": False,
                                "options": {"temperature": 0.4, "num_predict": 120,
                                            "stop": ["QUESTION:", "ANSWER:"]}},
                          timeout=60)
        r.raise_for_status()
        raw = r.json().get("response", "")
    except Exception:
        return []
    out = []
    for line in raw.splitlines():
        q = line.strip().lstrip("0123456789.-)*• ").strip().strip('"').strip()
        if len(q) >= 8:
            out.append(q if q.endswith("?") else q + "?")
        if len(out) == 3:
            break
    return out

@app.post("/followups")
def followups(req: FollowupReq):
    return {"questions": generate_followups(req.question, req.answer)}

# ── TMF630 conformance checker ───────────────────────────────────────────────────
@app.post("/conformance", tags=["Conformance"])
async def conformance_check(file: UploadFile = File(...)):
    """Upload an OpenAPI/Swagger spec → get a TMF630 conformance report (score + findings)."""
    raw  = await file.read()
    text = raw.decode("utf-8", errors="ignore")
    spec = None
    try:
        spec = json.loads(text)
    except Exception:
        try:
            import yaml
            spec = yaml.safe_load(text)
        except Exception:
            spec = None
    if not isinstance(spec, dict):
        raise HTTPException(400, "Could not parse that file as an OpenAPI/Swagger spec (JSON or YAML).")
    if not spec.get("paths") and not (spec.get("components") or spec.get("definitions")):
        raise HTTPException(400, "That doesn't look like an OpenAPI spec — no paths or schemas found.")
    report = conformance.check_spec(spec)
    report["filename"] = file.filename
    return report

class ConformanceTextReq(BaseModel):
    content: str
    filename: str = ""

@app.post("/conformance/text", tags=["Conformance"])
def conformance_check_text(req: ConformanceTextReq):
    """Same TMF630 check as /conformance but from a JSON body — used by the VS Code
    extension to validate the active editor's spec without a multipart upload."""
    text = req.content or ""
    if len(text) > 2_000_000:
        raise HTTPException(413, "Spec too large to check (over 2 MB).")
    spec = None
    try:
        spec = json.loads(text)
    except Exception:
        try:
            import yaml
            spec = yaml.safe_load(text)
        except Exception:
            spec = None
    if not isinstance(spec, dict):
        raise HTTPException(400, "Could not parse as an OpenAPI/Swagger spec (JSON or YAML).")
    if not spec.get("paths") and not (spec.get("components") or spec.get("definitions")):
        raise HTTPException(400, "Not an OpenAPI spec — no paths or schemas found.")
    report = conformance.check_spec(spec)
    report["filename"] = req.filename
    try:
        report["profile"] = tmf_profile.compare_to_canonical(spec)
    except Exception:
        report["profile"] = {"detected": None}
    return report


@app.post("/conformance/profile", tags=["Conformance"])
def conformance_profile(req: ConformanceTextReq):
    """Detect which TMF API a spec is and diff it against the canonical TMF spec
    (deterministic, no LLM). Powers the 'vs TMFxxx reference' coverage signal."""
    text = req.content or ""
    if len(text) > 2_000_000:
        raise HTTPException(413, "Spec too large to check (over 2 MB).")
    spec = None
    try:
        spec = json.loads(text)
    except Exception:
        try:
            import yaml
            spec = yaml.safe_load(text)
        except Exception:
            spec = None
    if not isinstance(spec, dict):
        raise HTTPException(400, "Could not parse as an OpenAPI/Swagger spec (JSON or YAML).")
    return tmf_profile.compare_to_canonical(spec)


class PortfolioReq(BaseModel):
    specs: list = []     # [{filename, content}, …]

@app.post("/conformance/portfolio", tags=["Conformance"])
def conformance_portfolio(req: PortfolioReq):
    """API estate X-ray: structural + profile conformance across many specs, rolled up
    into a board-ready report (markdown included). Deterministic, no LLM."""
    specs = (req.specs or [])[:300]
    report = xray.build_portfolio(specs)
    report["markdown"] = xray.render_markdown(report)
    return report

class ConformanceFixReq(BaseModel):
    content: str
    ids: Optional[list] = None     # specific rule ids to fix; None → fix all fixable

@app.post("/conformance/fix", tags=["Conformance"])
def conformance_fix(req: ConformanceFixReq):
    """Deterministically auto-fix TMF630 violations in a spec (no LLM) and return the
    corrected spec in its original format, plus the re-checked score."""
    content, fmt, spec = req.content or "", "json", None
    try:
        spec = json.loads(content)
    except Exception:
        try:
            import yaml
            spec = yaml.safe_load(content); fmt = "yaml"
        except Exception:
            spec = None
    if not isinstance(spec, dict):
        raise HTTPException(400, "Could not parse as an OpenAPI/Swagger spec (JSON or YAML).")
    fixed = conformance.fix_spec(spec, req.ids)
    report = conformance.check_spec(spec)
    if fmt == "yaml":
        import yaml
        out = yaml.safe_dump(spec, sort_keys=False, default_flow_style=False, allow_unicode=True, width=100)
    else:
        out = json.dumps(spec, indent=2, ensure_ascii=False)
    return {"content": out, "fixed": fixed, "format": fmt,
            "score": report["score"], "summary": report["summary"]}

@app.post("/conformance/scaffold", tags=["Conformance"])
def conformance_scaffold(req: ConformanceFixReq):
    """Complete a partial spec by merging the MISSING operations + resource attributes
    (and their schemas) from the detected canonical TMF spec — deterministic, no LLM.
    Returns the merged spec in its original format + what was added + coverage delta."""
    content = req.content or ""
    if len(content) > 2_000_000:
        raise HTTPException(413, "Spec too large to scaffold (over 2 MB).")
    fmt, spec = "json", None
    try:
        spec = json.loads(content)
    except Exception:
        try:
            import yaml
            spec = yaml.safe_load(content); fmt = "yaml"
        except Exception:
            spec = None
    if not isinstance(spec, dict):
        raise HTTPException(400, "Could not parse as an OpenAPI/Swagger spec (JSON or YAML).")
    res = tmf_profile.scaffold_from_canonical(spec)
    if not res.get("detected") or not res.get("spec"):
        return {"detected": None, "added": {"operations": [], "fields": [], "components": 0}}
    merged = res["spec"]
    if fmt == "yaml":
        import yaml
        out = yaml.safe_dump(merged, sort_keys=False, default_flow_style=False, allow_unicode=True, width=100)
    else:
        out = json.dumps(merged, indent=2, ensure_ascii=False)
    return {"content": out, "format": fmt, "detected": res["detected"], "added": res["added"],
            "coverage_before": res["coverage_before"], "coverage_after": res["coverage_after"]}

@app.post("/conformance/component", tags=["ODA"])
def conformance_component(req: ConformanceFixReq):
    """ODA Component conformance: parse an ODA Component manifest (.component.yaml) and
    score each TM Forum Open API it exposes/depends on. Deterministic, no LLM."""
    text = req.content or ""
    if len(text) > 2_000_000:
        raise HTTPException(413, "Manifest too large.")
    rep = oda.check_component(text)
    if not rep.get("component"):
        raise HTTPException(400, rep.get("error", "Not an ODA Component manifest (kind: Component)."))
    rep["markdown"] = oda.render_markdown(rep)
    return rep

@app.get("/oda/components", tags=["ODA"])
def oda_components():
    """The official ODA component map (TM Forum v1.0.0 — 35 components in 6 functional
    blocks), enriched with exposed/dependent APIs where reference manifests exist. Each
    component additionally carries DERIVED availability flags — specification_available,
    contract_available, supported_execution — computed from the authoritative backend
    (the adapter's execution gate + the vendored canonical specs), never hardcoded. This
    is a read-only catalog view; it does not touch the CTK execution flow or verdict."""
    cat = oda.catalog()
    supported = oda_ctk_jobs.oda_ctk_adapter.SUPPORTED_COMPONENTS
    components = [{**c, **oda.component_status(c, supported)} for c in cat.get("components", [])]
    return {**cat, "components": components}

# ── ODA Component CTK conformance (execution-backed — SEPARATE from schema conformance) ─
# A distinct conformance mode: resolve the canonical ODA Component contract, run the
# real TM Forum Component CTK against a DEPLOYED component, and normalize the CTK's own
# results deterministically. This does NOT touch conformance.py / tmf_profile.py /
# oda.check_component (the OpenAPI schema engine) — it is a different question ("does the
# deployed component pass its CTKs?"), not schema comparison.
class ODACTKConfig(BaseModel):
    company_name: str = ""
    product_name: str = ""
    product_url: str = ""
    product_version: str = ""
    headers: dict = {}                     # may include Authorization — redacted before persist/return
    payloads: dict = {}
    reject_unauthorized: bool = False

class ODACTKJobReq(BaseModel):
    component_id: str = "TMFC043"          # Phase 1: TMFC043 golden path only
    component_version: Optional[str] = None
    release_name: str = ""                 # deployed component's Helm release (required)
    namespace: str = "components"
    run_exposed_optional: bool = False
    run_dependent_optional: bool = False
    run_security_optional: bool = False
    ctkconfig: ODACTKConfig = ODACTKConfig()

def _ctk_req(req: "ODACTKJobReq") -> dict:
    return {
        "component_id": req.component_id, "component_version": req.component_version,
        "release_name": req.release_name, "namespace": req.namespace,
        "run_exposed_optional": req.run_exposed_optional,
        "run_dependent_optional": req.run_dependent_optional,
        "run_security_optional": req.run_security_optional,
        "ctkconfig": req.ctkconfig.dict(),
    }

@app.get("/oda/components/{component_id}/contract", tags=["ODA CTK"])
def oda_component_contract_endpoint(component_id: str):
    """Deterministic canonical ODA Component contract (mandatory/optional/dependent API
    requirements + events) from the vendored canonical YAML. No LLM, no RAG. Public/
    unauthenticated like its sibling GET /oda/components — deterministic, non-secret,
    stateless data (no job is created, nothing is stored, nothing to own)."""
    contract = oda_component_contract.resolve_contract(component_id)
    if contract.get("status") != "RESOLVED":
        reason = oda_ctk_sanitize.sanitize_reason(contract.get("reason", ""))
        raise HTTPException(404 if contract.get("status") == "UNSUPPORTED" else 422,
                            f"{contract.get('status')}: {reason}")
    return contract

# The 4 endpoints below are stateful and/or execute real infra tooling (subprocess,
# outbound calls using caller-supplied credentials) — unlike the stateless read above,
# these require login (Phase 4 audit: they previously had no auth boundary at all,
# inconsistent with the rest of the authenticated app surface). Job read endpoints are
# additionally bound to the creating user (or an admin) — see oda_ctk_jobs._owned().

@app.post("/oda/conformance/jobs/validate", tags=["ODA CTK"])
def oda_ctk_validate(req: ODACTKJobReq, user: dict = Depends(current_user)):
    """Adapter DRY RUN: resolve the contract, validate config, generate CHANGE_ME.json in
    an isolated workspace, verify framework files, list expected mandatory CTKs and
    prerequisites — WITHOUT invoking helm/kubectl/docker/npm or the CTK executor.
    READY_TO_EXECUTE means ready — NOT ODA-conformant."""
    return oda_ctk_sanitize.sanitize_dry_run(oda_ctk_jobs.oda_ctk_adapter.dry_run(_ctk_req(req)))

@app.post("/oda/conformance/jobs", tags=["ODA CTK"])
def oda_ctk_create_job(req: ODACTKJobReq, user: dict = Depends(current_user)):
    """Create a job and start the REAL CTK execution (background). Returns the job
    immediately; poll GET /oda/conformance/jobs/{id}. An execution/infra failure is
    classified honestly as FAILED_EXECUTION / EXECUTION_ERROR — never a conformance FAIL."""
    return oda_ctk_jobs.create_execution_job(_ctk_req(req), owner_id=user["id"])

@app.get("/oda/conformance/jobs/{job_id}", tags=["ODA CTK"])
def oda_ctk_get_job(job_id: str, user: dict = Depends(current_user)):
    job = oda_ctk_jobs.get_job(job_id, requester_id=user["id"], is_admin=(user.get("role") == "admin"))
    if not job:
        raise HTTPException(404, "job not found")
    return job

@app.get("/oda/conformance/jobs/{job_id}/results", tags=["ODA CTK"])
def oda_ctk_get_results(job_id: str, user: dict = Depends(current_user)):
    res = oda_ctk_jobs.get_results(job_id, requester_id=user["id"], is_admin=(user.get("role") == "admin"))
    if not res:
        raise HTTPException(404, "job not found")
    return res

# ── Documents ──────────────────────────────────────────────────────────────────
@app.get("/documents/library")
def documents_library():
    """
    Return the pre-indexed knowledge base documents grouped by source folder.
    These are the docs ingested from the data/ directory — read-only.
    """
    data_dir = Path("./data")
    upload_files = {d["file"] for d in load_uploads()}

    # Build filename → top-level folder mapping by scanning data/
    file_to_folder: dict[str, str] = {}
    if data_dir.exists():
        for folder in data_dir.iterdir():
            if folder.is_dir() and not folder.name.startswith("."):
                for f in folder.rglob("*"):
                    if f.is_file() and not f.name.startswith("."):
                        file_to_folder[f.name] = folder.name

    # Get chunk counts per file from ChromaDB (metadata segment, always accurate)
    col     = get_collection()
    records = col.scan()

    file_stats: dict[str, dict] = {}
    for meta in (r["metadata"] for r in records):
        fname = meta.get("file", "")
        if not fname or fname in upload_files or meta.get("origin"):
            continue                            # skip user-added sources (files/repos/web)
        if fname not in file_stats:
            folder = file_to_folder.get(fname, "other")
            # Human-readable name: strip extension + normalise
            display = fname
            for ext in (".swagger.json", ".oas.yaml", ".oas.json", ".yaml",
                        ".yml", ".json", ".md", ".txt"):
                if display.endswith(ext):
                    display = display[: -len(ext)]
                    break
            display = display.replace("-", " ").replace("_", " ")
            file_stats[fname] = {"file": fname, "name": display, "folder": folder,
                                 "category": _kb_category(meta.get("spec_id", ""), fname, folder),
                                 "chunks": 0, "status": "indexed"}
        file_stats[fname]["chunks"] += 1

    # Group by domain category (TM Forum / ODA / MEF / …) — same folders as user sources
    groups: dict[str, dict] = {}
    for stat in file_stats.values():
        g = stat["category"]
        if g not in groups:
            groups[g] = {"id": g, "name": g, "files": 0, "chunks": 0, "documents": []}
        groups[g]["files"]  += 1
        groups[g]["chunks"] += stat["chunks"]
        groups[g]["documents"].append(stat)

    CAT_ORDER = ["TM Forum", "ODA", "MEF", "ETSI", "3GPP", "IETF", "Other"]
    sorted_groups = sorted(groups.values(),
                           key=lambda x: CAT_ORDER.index(x["name"]) if x["name"] in CAT_ORDER else 99)
    for grp in sorted_groups:
        grp["documents"] = sorted(grp["documents"], key=lambda x: -x["chunks"])[:200]

    return {
        "total_files":  len(file_stats),
        "total_chunks": sum(s["chunks"] for s in file_stats.values()),
        "groups": sorted_groups,
    }

# ── Document preview (read the actual file content in-app) ───────────────────────
_KB_PATHS: dict[str, str] = {}      # basename → full path under data/ (lazy, cleared on refresh)

def _kb_path_index() -> dict[str, str]:
    global _KB_PATHS
    if not _KB_PATHS:
        data_dir = Path("./data")
        if data_dir.exists():
            for p in data_dir.rglob("*"):
                if p.is_file() and not p.name.startswith("."):
                    _KB_PATHS.setdefault(p.name, str(p))
    return _KB_PATHS

def _resolve_doc_file(file: str) -> Optional[Path]:
    """Map a document's name to a real file on disk — uploads first, then the KB.
    Only the basename is used, so a request can never escape uploads/ or data/."""
    name = Path(file or "").name
    if not name:
        return None
    up = UPLOADS_DIR / name                       # an uploaded original
    if up.is_file():
        return up
    kb = _kb_path_index().get(name)               # a bundled KB spec/doc
    if kb and Path(kb).is_file():
        return Path(kb)
    return None

PREVIEW_MAX_CHARS = 600_000                        # ~600 KB cap for inline text preview

@app.get("/documents/raw")
def document_raw(file: str):
    """Return a document's content as text for the in-app reader. Text files are
    returned verbatim (JSON pretty-printed); PDF/DOCX/XLSX/PPTX are returned as
    extracted text. Binary originals are streamed via /documents/file instead."""
    path = _resolve_doc_file(file)
    if not path:
        raise HTTPException(404, "Document not found")
    ext = path.suffix.lower()
    if   ext in (".md", ".markdown", ".rst"):           kind, lang = "markdown", "markdown"
    elif ext == ".json":                                 kind, lang = "code", "json"
    elif ext in (".yaml", ".yml"):                       kind, lang = "code", "yaml"
    elif ext == ".csv":                                  kind, lang = "code", "csv"
    elif ext in (".pdf", ".docx", ".xlsx", ".pptx"):     kind, lang = "extracted", ext.lstrip(".")
    else:                                                 kind, lang = "text", ""

    try:
        if kind == "extracted":
            content = extract_text(path.read_bytes(), path.name)
        else:
            content = path.read_text(encoding="utf-8", errors="ignore")
            if ext == ".json":                          # pretty-print for readability
                try:    content = json.dumps(json.loads(content), indent=2, ensure_ascii=False)
                except Exception: pass
    except Exception as e:
        raise HTTPException(500, f"Could not read document: {e}")

    truncated = len(content) > PREVIEW_MAX_CHARS
    return {"name": path.name, "ext": ext,
            "kind": "text" if kind == "extracted" else kind, "lang": lang,
            "extracted": kind == "extracted", "bytes": path.stat().st_size,
            "truncated": truncated, "content": content[:PREVIEW_MAX_CHARS],
            "downloadable": ext in (".pdf", ".docx", ".xlsx", ".pptx")}

@app.get("/documents/file")
def document_file(file: str, download: bool = False):
    """Serve a document's original bytes — used to embed PDFs inline and to download originals."""
    path = _resolve_doc_file(file)
    if not path:
        raise HTTPException(404, "Document not found")
    media = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    disp  = "attachment" if download else "inline"
    return FileResponse(str(path), media_type=media,
                        headers={"Content-Disposition": f'{disp}; filename="{path.name}"'})

def _confidence(chunks: list) -> dict:
    """How well-grounded an answer is, from retrieval distances (cosine; lower = better).
    Returns {level: high|medium|low, score: 0-100, strong: N} for a UI trust indicator."""
    if not chunks:
        return {"level": "low", "score": 0, "strong": 0}
    dists  = sorted(c.get("distance", 1.0) for c in chunks)
    best   = dists[0]
    strong = sum(1 for d in dists if d < 0.34)        # chunks that match strongly (cf. DOMAIN_GUARD_DIST)
    score  = max(0, min(100, round((0.60 - best) / 0.40 * 100)))   # 0.20→100, 0.60→0
    if   best < 0.33 and strong >= 2: level = "high"
    elif best < 0.43:                 level = "medium"
    else:                             level = "low"
    return {"level": level, "score": score, "strong": strong}

def _categorize(*hints) -> str:
    """Best-effort domain bucket for a USER-ADDED source (defaults to 'Other')."""
    blob = " ".join(h for h in hints if h).lower()
    if "tmforum" in blob or "tm forum" in blob or "tmf" in blob: return "TM Forum"
    if "oda" in blob:                    return "ODA"
    if "mef" in blob:                    return "MEF"
    if "etsi" in blob or "nfv" in blob:  return "ETSI"
    if "3gpp" in blob:                   return "3GPP"
    if "ietf" in blob or "rfc" in blob:  return "IETF"
    return "Other"

def _kb_category(spec_id: str, file: str, folder: str) -> str:
    """Domain bucket for a base Knowledge Base doc. Defaults to 'TM Forum' since
    the bundled KB is TM Forum-sourced; splits out ODA/MEF/etc. when detectable."""
    blob = f"{spec_id} {file} {folder}".lower()
    if "oda" in blob:                    return "ODA"
    if "mef" in blob:                    return "MEF"
    if "etsi" in blob or "nfv" in blob:  return "ETSI"
    if "3gpp" in blob:                   return "3GPP"
    if "ietf" in blob or "rfc" in blob:  return "IETF"
    return "TM Forum"

@app.get("/documents")
def list_documents():
    docs = load_uploads()
    # Overlay live in-memory status (more up-to-date than persisted JSON)
    for d in docs:
        fn = d.get("file","")
        if fn in PROCESSING_STATUS:
            live = PROCESSING_STATUS[fn]
            d["status"] = live["status"]
            d["chunks"] = live.get("chunks", d.get("chunks", 0))
        # Auto-organize: tag each source with a domain bucket (TM Forum / ODA / MEF / …)
        d["category"] = d.get("category") or _categorize(d.get("name"), d.get("url"), d.get("file"))
        d["folder"]   = d.get("folder") or ""     # user-assigned folder ("" = Uncategorized)
    return {"documents": docs}

# ── Folders ──────────────────────────────────────────────────────────────────────
class FolderReq(BaseModel):
    name: str

class MoveReq(BaseModel):
    folder: str = ""

@app.get("/folders")
def list_folders():
    docs = load_uploads()
    counts: dict = {}
    for d in docs:
        f = d.get("folder") or ""
        counts[f] = counts.get(f, 0) + 1
    return {"folders": [{"name": n, "count": counts.get(n, 0)} for n in load_folders()],
            "uncategorized": counts.get("", 0)}

@app.post("/folders")
def create_folder(req: FolderReq, _admin: dict = Depends(require_admin)):
    name = (req.name or "").strip()
    if not name:           raise HTTPException(400, "Folder name required")
    if len(name) > 60:     raise HTTPException(400, "Folder name too long (max 60)")
    folders = load_folders()
    if name not in folders:
        folders.append(name); save_folders(folders)
    return {"folders": load_folders()}

@app.delete("/folders/{name}")
def delete_folder(name: str, _admin: dict = Depends(require_admin)):
    save_folders([f for f in load_folders() if f != name])
    docs = load_uploads()                  # its documents fall back to Uncategorized (never deleted)
    for d in docs:
        if d.get("folder") == name:
            d["folder"] = ""
    save_uploads(docs)
    return {"ok": True, "folders": load_folders()}

@app.put("/folders/{name}")
def rename_folder(name: str, req: FolderReq, _admin: dict = Depends(require_admin)):
    """Rename a folder and move all of its documents to the new name."""
    new = (req.name or "").strip()
    if not new:                       raise HTTPException(400, "Folder name required")
    if len(new) > 60:                 raise HTTPException(400, "Folder name too long (max 60)")
    folders = load_folders()
    if name not in folders:           raise HTTPException(404, "Folder not found")
    if new == name:                   return {"ok": True, "folders": load_folders()}
    if new in folders:                raise HTTPException(409, "A folder with that name already exists")
    save_folders([new if f == name else f for f in folders])
    docs = load_uploads()             # carry the folder's documents over to the new name
    for d in docs:
        if d.get("folder") == name:
            d["folder"] = new
    save_uploads(docs)
    return {"ok": True, "folders": load_folders()}

@app.post("/documents/{file_name}/folder")
def move_document(file_name: str, req: MoveReq, _admin: dict = Depends(require_admin)):
    folder = (req.folder or "").strip()
    docs, found = load_uploads(), False
    for d in docs:
        if d.get("file") == file_name:
            d["folder"] = folder; found = True
    if not found:
        raise HTTPException(404, "Document not found")
    save_uploads(docs)
    if folder and folder not in load_folders():
        save_folders(load_folders() + [folder])
    return {"ok": True}

@app.post("/folders/{name}/refresh")
def refresh_folder(name: str, background_tasks: BackgroundTasks, _admin: dict = Depends(require_admin)):
    """Re-fetch every repo/web source in a folder (folder-level refresh)."""
    sources = [d for d in load_uploads()
               if (d.get("folder") or "") == name and d.get("type") in ("repo", "web") and d.get("url")]
    if not sources:
        return {"refreshed": 0, "note": "No repo or web sources in this folder."}
    recs = get_collection().scan()
    for u in sources:
        origin = u.get("origin") or u.get("file")
        ids = [r["id"] for r in recs if r["metadata"].get("origin") == origin]
        if ids:
            get_collection().delete(ids=ids)
        PROCESSING_STATUS[origin] = {"status": "processing", "chunks": 0}
        background_tasks.add_task(ingest_repo_bg if u["type"] == "repo" else ingest_web_bg,
                                  u["url"], origin, u.get("name") or origin)
    cache_clear()
    return {"refreshed": len(sources)}

@app.post("/documents/upload")
async def upload_document(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    content     = await file.read()
    file_name   = file.filename or "unnamed_file"
    source_name = file_name.rsplit(".", 1)[0].replace("_"," ").replace("-"," ")
    size_mb     = round(len(content) / 1_048_576, 1)

    # Persist raw bytes
    (UPLOADS_DIR / file_name).write_bytes(content)

    # Optimistically register as processing
    entry = {"file": file_name, "name": source_name,
             "size": f"{size_mb} MB", "status": "processing",
             "chunks": 0, "uploaded": time.strftime("%b %d, %Y")}
    upsert_upload(entry)
    PROCESSING_STATUS[file_name] = {"status": "processing", "chunks": 0}

    background_tasks.add_task(ingest_file_bg, content, file_name, source_name, size_mb)
    return {"file": file_name, "status": "processing"}

@app.get("/documents/{file_name}/status")
def document_status(file_name: str):
    if file_name in PROCESSING_STATUS:
        return PROCESSING_STATUS[file_name]
    for d in load_uploads():
        if d.get("file") == file_name:
            return {"status": d.get("status","unknown"), "chunks": d.get("chunks",0)}
    raise HTTPException(404, "Document not found")

@app.delete("/documents/{file_name}")
def delete_document(file_name: str):
    col       = get_collection()
    to_delete = [
        r["id"] for r in col.scan()
        if r["metadata"].get("origin") == file_name or r["metadata"].get("file") == file_name
    ]
    if to_delete:
        col.delete(ids=to_delete)
        # Keep the singleton — it already reflects the deletion in memory.

    remove_upload_record(file_name)
    PROCESSING_STATUS.pop(file_name, None)
    cache_clear()
    raw = UPLOADS_DIR / file_name
    if raw.exists(): raw.unlink()

    return {"deleted_chunks": len(to_delete), "file": file_name}

# ── Branding ───────────────────────────────────────────────────────────────────
@app.get("/branding")
def get_branding():
    if BRANDING_FILE.exists():
        try: return json.loads(BRANDING_FILE.read_text())
        except: pass
    return DEFAULT_BRANDING

@app.post("/branding")
def save_branding_endpoint(body: dict, _admin: dict = Depends(require_admin)):
    cleaned = {k: v for k, v in body.items()
               if k in {"companyName","tagline","primaryColor"} and isinstance(v, str)}
    merged  = {**DEFAULT_BRANDING, **cleaned}
    BRANDING_FILE.write_text(json.dumps(merged, indent=2))
    return merged

# ── Chat config (placeholder / suggestions / domain persona) ─────────────────
@app.get("/chat-config")
def get_chat_config():
    return load_chat_config()

@app.post("/chat-config")
def save_chat_config(body: dict, _admin: dict = Depends(require_admin)):
    cfg = load_chat_config()
    if isinstance(body.get("placeholder"), str): cfg["placeholder"] = body["placeholder"][:200]
    if isinstance(body.get("persona"), str):     cfg["persona"]     = body["persona"][:600]
    if isinstance(body.get("suggestions"), list):
        cfg["suggestions"] = [str(s).strip()[:200] for s in body["suggestions"] if str(s).strip()][:8]
    CHATCFG_FILE.write_text(json.dumps(cfg, indent=2))
    return cfg

# ── Conversations (per-user chat history) ────────────────────────────────────
def _convos_path(uid: str):
    return CONVOS_DIR / f"{re.sub(r'[^A-Za-z0-9_-]', '_', uid)}.json"

def _load_convos(uid: str) -> list:
    p = _convos_path(uid)
    if p.exists():
        try:    return json.loads(p.read_text())
        except: pass
    return []

def _save_convos(uid: str, convos: list):
    _convos_path(uid).write_text(json.dumps(convos, indent=2))

@app.get("/conversations")
def list_conversations(user: dict = Depends(current_user)):
    convos = _load_convos(user["id"])
    return {"conversations": sorted(
        [{"id": c["id"], "title": c.get("title", "New chat"),
          "updated": c.get("updated", 0), "count": len(c.get("messages", [])),
          "pinned": bool(c.get("pinned")), "project_id": c.get("project_id", "")}
         for c in convos],
        key=lambda x: (not x["pinned"], -x["updated"]))}   # pinned first, then most recent

@app.get("/conversations/search")
def search_conversations(q: str = "", user: dict = Depends(current_user)):
    """Full-text search across a user's chat titles AND message contents, with a snippet
    of the first match. Defined before /{cid} so 'search' isn't read as an id."""
    ql = (q or "").strip().lower()
    if len(ql) < 2:
        return {"results": []}
    out = []
    for c in _load_convos(user["id"]):
        title = c.get("title", "") or "Untitled"
        snippet, matched = "", ql in title.lower()
        for m in c.get("messages", []):
            content = m.get("content") or ""
            idx = content.lower().find(ql)
            if idx >= 0:
                matched = True
                start = max(0, idx - 40)
                snippet = (("…" if start > 0 else "") + content[start:idx + len(ql) + 70]).replace("\n", " ").strip()
                break
        if matched:
            out.append({"id": c["id"], "title": title, "updated": c.get("updated", 0),
                        "pinned": bool(c.get("pinned")), "project_id": c.get("project_id", ""),
                        "snippet": snippet})
    out.sort(key=lambda x: -x["updated"])
    return {"results": out[:50]}

@app.get("/conversations/{cid}")
def get_conversation(cid: str, user: dict = Depends(current_user)):
    for c in _load_convos(user["id"]):
        if c["id"] == cid:
            return c
    raise HTTPException(404, "Conversation not found")

@app.post("/conversations")
def create_conversation(body: dict, user: dict = Depends(current_user)):
    convos = _load_convos(user["id"])
    c = {"id": uuid.uuid4().hex[:12],
         "title": (body.get("title") or "New chat")[:80],
         "messages": body.get("messages") or [],
         "project_id": body.get("project_id") or "",
         "updated": int(time.time())}
    convos.append(c)
    _save_convos(user["id"], convos)
    return c

@app.put("/conversations/{cid}")
def update_conversation(cid: str, body: dict, user: dict = Depends(current_user)):
    convos = _load_convos(user["id"])
    for c in convos:
        if c["id"] == cid:
            changed_msgs = isinstance(body.get("messages"), list)
            if changed_msgs:
                c["messages"] = body["messages"][:200]
            if body.get("rename"):                                  # explicit manual rename — sticks
                c["title"] = (str(body.get("title") or "").strip() or "Untitled")[:80]
                c["custom_title"] = True
            elif body.get("title") and not c.get("custom_title"):    # auto title from 1st message
                c["title"] = str(body["title"])[:80]
            if "pinned" in body:
                c["pinned"] = bool(body["pinned"])
            if "project_id" in body:                                # assign / remove from a project
                c["project_id"] = body.get("project_id") or ""
            if changed_msgs:                                         # only real activity bumps recency
                c["updated"] = int(time.time())
            _save_convos(user["id"], convos)
            return {"id": c["id"], "title": c["title"], "pinned": bool(c.get("pinned")),
                    "project_id": c.get("project_id", ""),
                    "updated": c.get("updated", 0), "count": len(c.get("messages", []))}
    raise HTTPException(404, "Conversation not found")

@app.delete("/conversations/{cid}")
def delete_conversation(cid: str, user: dict = Depends(current_user)):
    convos = _load_convos(user["id"])
    remaining = [c for c in convos if c["id"] != cid]
    if len(remaining) == len(convos):
        raise HTTPException(404, "Conversation not found")
    _save_convos(user["id"], remaining)
    return {"ok": True}

# ── Projects (group chats + per-project standing notes / memory) ─────────────────
def _projects_path(uid: str):
    return PROJECTS_DIR / f"{re.sub(r'[^A-Za-z0-9_-]', '_', uid)}.json"

def _load_projects(uid: str) -> list:
    p = _projects_path(uid)
    if p.exists():
        try:    return json.loads(p.read_text())
        except: pass
    return []

def _save_projects(uid: str, projects: list):
    _projects_path(uid).write_text(json.dumps(projects, indent=2))

@app.get("/projects")
def list_projects(user: dict = Depends(current_user)):
    projects = _load_projects(user["id"])
    counts: dict = {}
    for c in _load_convos(user["id"]):
        pid = c.get("project_id") or ""
        if pid:
            counts[pid] = counts.get(pid, 0) + 1
    return {"projects": sorted(
        [{"id": p["id"], "name": p.get("name", "Untitled project"),
          "instructions": p.get("instructions", ""), "updated": p.get("updated", 0),
          "count": counts.get(p["id"], 0)} for p in projects],
        key=lambda x: -x["updated"])}

@app.post("/projects")
def create_project(body: dict, user: dict = Depends(current_user)):
    name = (body.get("name") or "").strip()
    if not name:        raise HTTPException(400, "Project name required")
    if len(name) > 80:  raise HTTPException(400, "Project name too long (max 80)")
    projects = _load_projects(user["id"])
    p = {"id": uuid.uuid4().hex[:12], "name": name,
         "instructions": (body.get("instructions") or "")[:4000], "updated": int(time.time())}
    projects.append(p)
    _save_projects(user["id"], projects)
    return p

@app.put("/projects/{pid}")
def update_project(pid: str, body: dict, user: dict = Depends(current_user)):
    projects = _load_projects(user["id"])
    for p in projects:
        if p["id"] == pid:
            if body.get("name"):
                p["name"] = str(body["name"]).strip()[:80]
            if "instructions" in body:
                p["instructions"] = str(body.get("instructions") or "")[:4000]
            p["updated"] = int(time.time())
            _save_projects(user["id"], projects)
            return p
    raise HTTPException(404, "Project not found")

@app.delete("/projects/{pid}")
def delete_project(pid: str, user: dict = Depends(current_user)):
    projects = _load_projects(user["id"])
    if not any(p["id"] == pid for p in projects):
        raise HTTPException(404, "Project not found")
    _save_projects(user["id"], [p for p in projects if p["id"] != pid])
    convos = _load_convos(user["id"])              # its chats become unfiled — never deleted
    changed = False
    for c in convos:
        if c.get("project_id") == pid:
            c["project_id"] = ""; changed = True
    if changed:
        _save_convos(user["id"], convos)
    return {"ok": True}

# ── Users ──────────────────────────────────────────────────────────────────────
_DEFAULTS_BY_EMAIL = {u["email"]: u for u in DEFAULT_USERS}

def _ensure_passwords(users: list[dict]) -> bool:
    """Migrate stored users: give everyone a password hash (seed accounts use
    their known demo password) and backfill title/department from defaults."""
    changed = False
    for u in users:
        em = (u.get("email") or "").lower()
        if not u.get("password_hash"):
            pw = SEED_PASSWORDS.get(em, secrets.token_hex(8))
            u["password_hash"], u["salt"] = hash_password(pw)
            changed = True
        d = _DEFAULTS_BY_EMAIL.get(em)
        if d:
            for k in ("title", "department"):
                if not u.get(k) and d.get(k):
                    u[k] = d[k]; changed = True
    return changed

def _load_users() -> list[dict]:
    if USERS_FILE.exists():
        try:
            users = json.loads(USERS_FILE.read_text())
            if _ensure_passwords(users):
                _save_users(users)
            return users
        except: pass
    users = [dict(u) for u in DEFAULT_USERS]
    _ensure_passwords(users)
    _save_users(users)
    return users

def _save_users(users: list[dict]):
    USERS_FILE.write_text(json.dumps(users, indent=2))

# ── Auth endpoints ───────────────────────────────────────────────────────────
@app.post("/auth/login")
def auth_login(body: LoginReq):
    email = (body.email or "").lower().strip()
    users = _load_users()
    u = next((x for x in users if (x.get("email") or "").lower() == email), None)
    if not u or not verify_password(body.password, u.get("password_hash", ""), u.get("salt", "")):
        raise HTTPException(401, "Invalid email or password.")
    u["lastActive"] = "Now"
    _save_users(users)
    return {"token": make_token(u), "user": _public_user(u)}

@app.get("/auth/me")
def auth_me(user: dict = Depends(current_user)):
    return {"user": _public_user(user)}

@app.put("/auth/profile")
def auth_update_profile(body: dict, user: dict = Depends(current_user)):
    """Let a signed-in user edit their own name / title / department."""
    users = _load_users()
    for u in users:
        if u["id"] == user["id"]:
            for k in ("name", "title", "department"):
                if k in body and isinstance(body[k], str):
                    u[k] = body[k]
            _save_users(users)
            return {"user": _public_user(u)}
    raise HTTPException(404, "User not found")

@app.post("/auth/password")
def auth_change_password(body: dict, user: dict = Depends(current_user)):
    """Change the signed-in user's own password (verifies the current one)."""
    new = body.get("new_password", "")
    if len(new) < 6:
        raise HTTPException(400, "New password must be at least 6 characters.")
    users = _load_users()
    for u in users:
        if u["id"] == user["id"]:
            if not verify_password(body.get("current_password", ""), u.get("password_hash", ""), u.get("salt", "")):
                raise HTTPException(401, "Current password is incorrect.")
            u["password_hash"], u["salt"] = hash_password(new)
            _save_users(users)
            return {"ok": True}
    raise HTTPException(404, "User not found")

@app.get("/users")
def get_users():
    return {"users": [_public_user(u) for u in _load_users()]}

@app.post("/users")
def add_user(body: dict, _admin: dict = Depends(require_admin)):
    users = _load_users()
    ph, salt = hash_password(body.get("password") or secrets.token_hex(8))
    new_u = {"id": str(uuid.uuid4()), "name": body.get("name", ""),
             "email": body.get("email", ""), "role": body.get("role", "viewer"),
             "status": "inactive", "lastActive": "Never",
             "title": body.get("title", ""), "department": body.get("department", ""),
             "password_hash": ph, "salt": salt}
    users.append(new_u)
    _save_users(users)
    return _public_user(new_u)

@app.put("/users/{user_id}")
def update_user(user_id: str, body: dict, _admin: dict = Depends(require_admin)):
    users = _load_users()
    for u in users:
        if u["id"] == user_id:
            for k in ("name", "role", "status", "title", "department"):
                if k in body: u[k] = body[k]
            if body.get("password"):
                u["password_hash"], u["salt"] = hash_password(body["password"])
            _save_users(users)
            return _public_user(u)
    raise HTTPException(404, "User not found")

@app.delete("/users/{user_id}")
def delete_user(user_id: str, _admin: dict = Depends(require_admin)):
    users     = _load_users()
    remaining = [u for u in users if u["id"] != user_id]
    if len(remaining) == len(users):
        raise HTTPException(404, "User not found")
    _save_users(remaining)
    return {"ok": True}

# ── Query analytics ────────────────────────────────────────────────────────────
@app.get("/stats/queries")
def query_stats():
    return {"hits": QUERY_HITS}

@app.get("/stats/documents")
def stats_documents(period: str = "30d"):
    """Per-document query volume + average retrieval relevance over a real time window
    (7d/30d/90d), from the persisted dated store. Powers the Document Performance tab."""
    n = {"7d": 7, "30d": 30, "90d": 90}.get(period, 30)
    cutoff = set(time.strftime("%Y-%m-%d", time.localtime(time.time() - i * 86400)) for i in range(n))
    try:
        data = json.loads(DOC_PERF_FILE.read_text()) if DOC_PERF_FILE.exists() else {}
    except Exception:
        data = {}
    agg: dict = {}
    for day, bucket in data.items():
        if day not in cutoff:
            continue
        for name, e in bucket.items():
            a = agg.get(name) or {"queries": 0, "rel_sum": 0.0}
            a["queries"] += e.get("hits", 0)
            a["rel_sum"] += e.get("rel_sum", 0.0)
            agg[name] = a
    docs = {name: {"queries": a["queries"],
                   "relevance": round(a["rel_sum"] / a["queries"], 3) if a["queries"] else 0.0}
            for name, a in agg.items()}
    return {"period": period, "documents": docs}

# ── Feedback + analytics + KB refresh ────────────────────────────────────────
@app.post("/feedback")
def submit_feedback(body: dict, user: dict = Depends(current_user)):
    rating = body.get("rating")
    if rating not in ("up", "down"):
        raise HTTPException(400, "rating must be 'up' or 'down'")
    try:    fb = json.loads(FEEDBACK_FILE.read_text()) if FEEDBACK_FILE.exists() else []
    except Exception: fb = []
    fb.insert(0, {"ts": int(time.time()), "user": user.get("email"), "rating": rating,
                  "question": str(body.get("question", ""))[:200],
                  "sources": (body.get("sources") or [])[:5]})
    FEEDBACK_FILE.write_text(json.dumps(fb[:500], indent=2))
    return {"ok": True}

@app.get("/analytics")
def get_analytics(_admin: dict = Depends(require_admin)):
    try:    a = json.loads(ANALYTICS_FILE.read_text()) if ANALYTICS_FILE.exists() else {}
    except Exception: a = {}
    try:    fb = json.loads(FEEDBACK_FILE.read_text()) if FEEDBACK_FILE.exists() else []
    except Exception: fb = []
    top = sorted(QUERY_HITS.items(), key=lambda x: -x[1])[:10]

    # ── Adoption series ──────────────────────────────────────────────────────
    days_raw = a.get("days", {})
    today = time.time()
    last14 = [time.strftime("%Y-%m-%d", time.localtime(today - i * 86400)) for i in range(13, -1, -1)]
    series = [{"day": d[5:],                       # "MM-DD" for compact axis labels
               "q": days_raw.get(d, {}).get("q", 0),
               "a": days_raw.get(d, {}).get("a", 0)} for d in last14]
    last7 = set(time.strftime("%Y-%m-%d", time.localtime(today - i * 86400)) for i in range(7))
    wau_users = set()
    for d in last7:
        wau_users.update(u for u in days_raw.get(d, {}).get("users", {}) if u != "anonymous")
    last30 = set(time.strftime("%Y-%m-%d", time.localtime(today - i * 86400)) for i in range(30))
    per_user: dict = {}
    for d in last30:
        for u, n in days_raw.get(d, {}).get("users", {}).items():
            per_user[u] = per_user.get(u, 0) + n
    active = sorted(per_user.items(), key=lambda kv: -kv[1])[:10]

    return {
        "total_queries": a.get("total", 0),
        "answered":      a.get("answered", 0),
        "gaps":          a.get("gaps", [])[:20],
        "top_sources":   [{"name": k, "hits": v} for k, v in top],
        "series":        series,
        "wau":           len(wau_users),
        "active_users":  [{"user": u, "queries": n} for u, n in active],
        "feedback": {
            "up":   sum(1 for f in fb if f.get("rating") == "up"),
            "down": sum(1 for f in fb if f.get("rating") == "down"),
            "recent": fb[:10],
        },
    }

def _kb_refresh_job(mode: str = "selective", domain: Optional[str] = None):
    """Pull upstream + update the index. Runs synchronously (callers thread it).
    Shared by the manual /admin/refresh-kb endpoint and the scheduled auto-refresh."""
    global _spec_map, _spec_names
    PROCESSING_STATUS["__kb_refresh__"] = {"status": "processing", "mode": mode, "domain": domain or "all"}
    try:
        url_map = ingest.download_data(update=True)         # pull upstream changes
        if mode == "full" and not domain:
            result = {"chunks": ingest.build_index(url_map)}
        else:
            result = ingest.sync_index(url_map, domain=domain or None)   # only changed files
        with _chroma_lock:                     # force re-open of the updated collection
            reset_store()
            _spec_map = None; _spec_names = []
        _KB_PATHS.clear()                      # KB files changed → rebuild preview path index lazily
        tmf_profile.clear_index()              # canonical specs may have changed → rebuild the profile index
        cache_clear()
        PROCESSING_STATUS["__kb_refresh__"] = {"status": "indexed", "mode": mode, "domain": domain or "all", **result}
        return result
    except Exception as e:
        PROCESSING_STATUS["__kb_refresh__"] = {"status": "failed", "error": str(e)[:200]}
        raise

@app.post("/admin/refresh-kb")
def refresh_kb(mode: str = "selective", domain: str = "", _admin: dict = Depends(require_admin)):
    """Re-pull the TM Forum source repos and update the index.
    mode=selective (default): pull upstream + re-embed ONLY changed/removed files — fast.
    mode=full: drop and rebuild the whole index from scratch.
    domain (e.g. "TM Forum", "ODA"): scope a selective refresh to one KB folder only."""
    mode = "full" if (mode or "").lower() == "full" else "selective"
    domain = (domain or "").strip()
    if domain:                                  # per-folder refresh is always selective
        mode = "selective"
    if PROCESSING_STATUS.get("__kb_refresh__", {}).get("status") == "processing":
        return {"status": "processing"}
    threading.Thread(target=lambda: _kb_refresh_job(mode, domain or None), daemon=True).start()
    return {"status": "processing", "mode": mode, "domain": domain or "all"}

class RefreshSourceReq(BaseModel):
    origin: str

@app.post("/sources/refresh")
def refresh_source(req: RefreshSourceReq, background_tasks: BackgroundTasks, _admin: dict = Depends(require_admin)):
    """Re-fetch and re-index a single repo or web source (e.g. after it changed upstream)."""
    origin = (req.origin or "").strip()
    u = next((x for x in load_uploads() if (x.get("origin") or x.get("file")) == origin), None)
    if not u:
        raise HTTPException(404, "Source not found")
    typ, url, label = u.get("type"), u.get("url"), u.get("name") or origin
    if typ not in ("repo", "web") or not url:
        raise HTTPException(400, "Only Git-repo and web-link sources can be refreshed. Re-upload a file to update it.")
    # Drop the source's existing chunks first so removed files don't linger.
    col = get_collection()
    old_ids = [r["id"] for r in col.scan() if r["metadata"].get("origin") == origin]
    if old_ids:
        col.delete(ids=old_ids)
    cache_clear()
    PROCESSING_STATUS[origin] = {"status": "processing", "chunks": 0}
    background_tasks.add_task(ingest_repo_bg if typ == "repo" else ingest_web_bg, url, origin, label)
    return {"origin": origin, "status": "processing"}

@app.get("/admin/refresh-kb/status")
def refresh_kb_status(_admin: dict = Depends(require_admin)):
    return PROCESSING_STATUS.get("__kb_refresh__", {"status": "idle", "chunks": 0})

@app.post("/admin/cache/clear")
def admin_clear_cache(_admin: dict = Depends(require_admin)):
    """Development-safe response-cache reset: drops only cached /query and
    /query/stream answers (answer_cache.json). Use this to force every question
    through the real, current pipeline again — e.g. after a routing/grounding change,
    before trusting a live smoke test. Does NOT touch ChromaDB, indexed documents, or
    conversation history; use /admin/refresh-kb to reindex content instead."""
    cache_clear()
    return {"status": "cleared"}

# ── Scheduled auto-refresh ───────────────────────────────────────────────────────
SCHED_FILE = STORAGE_DIR / "refresh_schedule.json"

def _load_sched() -> dict:
    base = {"enabled": False, "interval": "weekly", "hour": 3, "last_run": 0, "next_run": 0}
    if SCHED_FILE.exists():
        try:    base.update(json.loads(SCHED_FILE.read_text()))
        except Exception: pass
    return base

def _save_sched(s: dict):
    try:    SCHED_FILE.write_text(json.dumps(s, indent=2))
    except Exception: pass

def _next_run_ts(s: dict) -> int:
    import datetime
    now  = datetime.datetime.now()
    hour = max(0, min(23, int(s.get("hour", 3))))
    target = now.replace(hour=hour, minute=0, second=0, microsecond=0)
    if s.get("interval") == "weekly":
        days = (6 - now.weekday()) % 7          # 6 = Sunday
        target += datetime.timedelta(days=days)
        if target <= now: target += datetime.timedelta(days=7)
    else:                                        # daily
        if target <= now: target += datetime.timedelta(days=1)
    return int(target.timestamp())

def _scheduler_loop():
    while True:
        try:
            s = _load_sched()
            if s.get("enabled") and s.get("next_run", 0) and time.time() >= s["next_run"]:
                if PROCESSING_STATUS.get("__kb_refresh__", {}).get("status") != "processing":
                    print("[scheduler] running scheduled KB refresh")
                    try:    _kb_refresh_job("selective", None)
                    except Exception as e: print(f"[scheduler] refresh failed: {e}")
                    s = _load_sched()
                    s["last_run"] = int(time.time())
                    s["next_run"] = _next_run_ts(s)
                    _save_sched(s)
        except Exception as e:
            print(f"[scheduler] loop error: {e}")
        time.sleep(300)                          # check every 5 minutes

class SchedReq(BaseModel):
    enabled: bool = False
    interval: str = "weekly"
    hour: int = 3

@app.get("/admin/refresh-schedule")
def get_refresh_schedule(_admin: dict = Depends(require_admin)):
    return _load_sched()

@app.put("/admin/refresh-schedule")
def set_refresh_schedule(req: SchedReq, _admin: dict = Depends(require_admin)):
    s = _load_sched()
    s["enabled"]  = bool(req.enabled)
    s["interval"] = "daily" if req.interval == "daily" else "weekly"
    s["hour"]     = max(0, min(23, int(req.hour)))
    s["next_run"] = _next_run_ts(s) if s["enabled"] else 0
    _save_sched(s)
    return s
